"""UX V11 — deck/PPTX export quality (spec/ux-contract.md §9): the synthesis→slides mapping uses
the master template's layout vocabulary (statement verdict, takeaway cards, chart/quote/voices
slides, numbered recommendation cards) and NO slide carries a wall of continuous prose — leads are
clamped renderer-side, markdown artifacts and bare artifact ids never reach a slide."""
import io

from pptx import Presentation

from sonaloop import services
from sonaloop.services._synthesis import _SYNTHESIS_EXPORT_LABELS
from sonaloop.services._synthesis_pptx import (
    _analytic_slides, _clamp_prose, _label_segments, _split_card,
)

# A showcase-shaped convergence synthesis: caps-label exec prose, stanced voices with verbatim
# council quotes, scored recommendations and label-led findings — the report shape the owner's
# screenshots came from.
_BIG = (
    "THE WINNING CONCEPT: A lunch window that stays visible in the day and slides audibly instead "
    "of dissolving in silence, with hard walls for fixed commitments. When the window hits a wall, "
    "the system says so honestly instead of sliding an illusion. "
    "WHY IT WINS: It makes the break as loud as its disruptors and productises the only defence "
    "observed in the field. The mid-fi validation was strict and every objection was retested. "
    "SURVEY BACKING (5 answers): The problem is real — 4/5 write off their break before 10 am on "
    "busy days. The visible window holds: 1 for, 3 conditional, 1 skeptical."
)


def _showcase(store):
    council = {
        "id": "c1", "created_at": "2026-06-10T00:00:00+00:00", "prompt": "Final check",
        "votes": [{"persona_id": "p1", "vote": "support", "reason": "works"},
                  {"persona_id": "p2", "vote": "conditional", "reason": "channel missing"},
                  {"persona_id": "p3", "vote": "skeptical", "reason": "wallboard decides"}],
        "statements": [
            {"persona_id": "p1", "text": "It calculates live.", "stance": {"value": 2, "label": "support"}},
            {"persona_id": "p2", "text": "Slack is missing.", "stance": {"value": 1, "label": "conditional"}},
            {"persona_id": "p3", "text": "Not my lever.", "stance": {"value": -1, "label": "skeptical"}},
        ],
    }
    store.insert_council_session(council)
    return services.record_synthesis(
        "Final solution presentation", "hmw", ["c1"],
        {
            "gesamtbild": _BIG,
            "positionierung": ("For desk workers with calendar autonomy whose lunch is the only "
                               "meeting without an advocate: the window advocate makes the lunch "
                               "window as loud as its disruptors, slides visibly instead of dying "
                               "silently, computes counter-proposals from the real day and belongs "
                               "to the person alone — never to the reporting chain above, not a "
                               "wellness ideal, not a KPI, a negotiation tool for the one meeting "
                               "without a lobby, which is why the team rule beats every individual "
                               "justification and the care framing is never acceptable without comment."),
            "findings": [
                {"kind": "pain_solver",
                 "text": "DYNAMIC COUNTER-PROPOSAL: productises the only successful defence observed "
                         "in the field (a 13:45 counter-proposal accepted without comment)."},
                {"kind": "recommendation",
                 "text": "BUILD-SPEC CRITERION 1 — Calendar & channel binding: disruptions must come "
                         "from calendar, inbox and Slack (Fabian protosession_14a1aa34d5562f24, "
                         "verified live).",
                 "score": {"effort": 4, "value": 5}},
                {"kind": "recommendation",
                 "text": "CRITERION 2 — Honest failure: when the window hits a hard wall at 15:00 the "
                         "day is marked as lost, no auto-slide.",
                 "score": {"effort": 2, "value": 4}},
                {"kind": "open_question",
                 "text": "SURVEILLANCE TIPPING POINT (biggest risk): all five name the same boundary — "
                         "the moment window data becomes a KPI, acceptance flips to rejection."},
            ],
            "statements": [
                {"persona_id": "p1",
                 "text": "The counter-proposal is **provably** computed live from my current window — "
                         "my `lo-fi` stress test, passed this time.",
                 "stance": {"value": 2, "label": "support"},
                 "refs": [{"kind": "council", "id": "c1",
                           "quote": "My break rarely tips with a bang, it gets deprioritised."}],
                 "meta": {"persona_name": "Fabian Drees", "segment": "Calendar-autonomous desk workers"}},
                {"persona_id": "p3",
                 "text": "The wallboard decides, not me — an individual advocate changes nothing about "
                         "two outages in the early shift.",
                 "stance": {"value": -1, "label": "skeptical"},
                 "refs": [],
                 "meta": {"persona_name": "Janine Wolf", "segment": "Shift operations (non-target)"}},
            ],
        },
        store=store)


def _slide_model(syn, store):
    return _analytic_slides(syn, store, _SYNTHESIS_EXPORT_LABELS["en"], False,
                            "Final solution presentation", "Report")


def _frame_texts(prs):
    """[(slide_index, frame_text), …] over every text frame of the rendered deck."""
    out = []
    for i, slide in enumerate(prs.slides):
        for sh in slide.shapes:
            if sh.has_text_frame:
                out.append((i, "\n".join(p.text for p in sh.text_frame.paragraphs)))
    return out


# ---------------------------------------------------------------- slide model


def test_deck_uses_master_layout_vocabulary(store):
    """Cover → verdict STATEMENT → takeaway cards → sentiment/stance CHARTS → quote + voices →
    finding cards → numbered recommendation cards + effort·value map: the report renders through
    the template vocabulary, not as content/prose slides."""
    syn = _showcase(store)
    kinds = [s["kind"] for s in _slide_model(syn, store)]
    assert kinds[0] == "cover"
    for kind in ("insight", "summary", "charts", "quote", "voices", "chart"):
        assert kind in kinds, f"missing {kind!r} slide (got {kinds})"


def test_verdict_statement_is_clamped_large_type(store):
    syn = _showcase(store)
    verdict = next(s for s in _slide_model(syn, store) if s["kind"] == "insight")
    assert verdict["eyebrow"] == "Verdict"
    assert 0 < len(verdict["statement"]) <= 240          # one sentence, large type — never a wall
    assert "THE WINNING CONCEPT:" not in verdict["statement"]   # the caps label is structure, not copy


def test_exec_summary_takeaways_from_authored_labels(store):
    """The caps-label spine of the authored exec prose becomes the takeaway cards (renderer-side
    splitting — headline + takeaways instead of a paragraph dump)."""
    syn = _showcase(store)
    summary = next(s for s in _slide_model(syn, store) if s["kind"] == "summary"
                   and s["heading"] == "Executive summary")
    titles = [it["title"] for it in summary["items"]]
    assert "WHY IT WINS" in titles and "SURVEY BACKING (5 answers)" in titles
    assert all(len(it["text"]) <= 260 for it in summary["items"])


def test_sentiment_and_stance_render_as_chart_slides(store):
    """Votes + contribution stances chart with scale-ordered, zero-free categories."""
    syn = _showcase(store)
    charts = next(s for s in _slide_model(syn, store) if s["kind"] == "charts")
    pie, bar = (it["chart"] for it in charts["items"])
    assert pie["type"] == "pie" and bar["type"] == "bar"
    assert pie["categories"] == ["Support", "Conditional", "Skeptical"]   # no zero buckets (V3)
    assert pie["values"] == [1, 1, 1] and bar["values"] == [1, 1, 1]


def test_voices_are_quote_slides_two_per_slide(store):
    syn = _showcase(store)
    slides = _slide_model(syn, store)
    quote = next(s for s in slides if s["kind"] == "quote")
    assert quote["attribution"] == "Fabian Drees"
    assert "deprioritised" in quote["text"]
    voices = [s for s in slides if s["kind"] == "voices"]
    assert voices and all(len(s["items"]) <= 2 for s in voices)
    fabian = voices[0]["items"][0]
    assert fabian["sentiment"] == "support" and fabian["sentiment_label"] == "Support"


def test_recommendations_are_numbered_cards_with_quiet_meta(store):
    syn = _showcase(store)
    rec_cards = next(s for s in _slide_model(syn, store)
                     if s["kind"] == "summary" and s.get("heading") == "Recommendations")
    assert rec_cards["items"][0]["title"].startswith("01 · ")
    assert rec_cards["items"][0]["meta"] == "Effort 4/5 · Value 5/5"
    # the effort·value map keeps SHORT legend labels (the heads), never the full prose
    rec_map = next(s for s in _slide_model(syn, store) if s["kind"] == "chart")
    assert all(len(p["label"]) <= 90 for p in rec_map["chart"]["points"])


# ------------------------------------------------------------- rendered deck


def test_no_rendered_frame_exceeds_the_prose_budget(store):
    """The §9 bar: no slide body is a text wall — every text frame stays within ~6 lines of
    continuous prose (≤700 chars per frame, ≤500 per paragraph)."""
    syn = _showcase(store)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx(syn["id"], store=store)))
    for i, text in _frame_texts(prs):
        assert len(text) <= 700, f"slide {i + 1}: frame carries {len(text)} chars of prose:\n{text[:200]}"
        for para in text.split("\n"):
            assert len(para) <= 500, f"slide {i + 1}: paragraph runs {len(para)} chars"


def test_no_markdown_artifacts_or_bare_ids_on_slides(store):
    syn = _showcase(store)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx(syn["id"], store=store)))
    deck_text = " ".join(t for _, t in _frame_texts(prs))
    assert "**" not in deck_text and "`" not in deck_text
    assert "protosession_" not in deck_text          # terminal ids are not slide copy
    assert "provably" in deck_text                   # the de-markdowned content survived


def test_project_section_prose_is_budgeted_with_report_footnote(store):
    """Project-report sections clamp to the slide budget and point to the full report instead of
    dumping the section body."""
    long_para = "The full analysis goes on and on with detail after detail. " * 40
    rep = {"id": "rlong", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-10T00:00:00+00:00", "lead": "", "council_ids": [],
           "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": "s1", "heading": "Findings", "markdown": long_para,
                         "citations": [], "source_study_ids": [], "figures": []}]}
    store.upsert_synthesis(rep)
    prs = Presentation(io.BytesIO(services.export_synthesis_pptx("rlong", store=store)))
    texts = [t for _, t in _frame_texts(prs)]
    assert all(len(t) <= 700 for t in texts)
    assert any("Details in the full report" in t for t in texts)


# ------------------------------------------------------------------- helpers


def test_split_card_never_breaks_words_or_times():
    head, body = _split_card("Kalender-autonome Schreibtisch-Arbeiter:innen (Zielsegment) — Stance: win")
    assert head == "Kalender-autonome Schreibtisch-Arbeiter:innen (Zielsegment)"
    assert body == "Stance: win"
    head, body = _split_card("Hard walls at 15:00 stay fixed: the window yields.")
    assert head == "Hard walls at 15:00 stay fixed" and body == "the window yields."
    assert _split_card("No separator here at all") == ("", "No separator here at all")


def test_clamp_prose_cuts_at_sentence_and_clause_boundaries():
    text = "First point stands. Second point follows. Third point would overflow the budget."
    out, truncated = _clamp_prose(text, 45)
    assert out == "First point stands. Second point follows." and truncated
    out, truncated = _clamp_prose("One enormous clause, with a second clause, and a third clause", 40)
    assert out.endswith("…") and len(out) <= 45 and truncated


def test_label_segments_finds_the_caps_spine():
    segs = _label_segments(_BIG)
    assert [s[0] for s in segs] == ["THE WINNING CONCEPT", "WHY IT WINS", "SURVEY BACKING (5 answers)"]
