"""The deck master template (vendored _deck.py ← sonaloop-design/deck.data.mjs) — every layout
renders to native, editable PPTX shapes, and the placeholder SAMPLE_SLIDES are the demo deck
the harness exposes via `sonaloop template-deck` / `make template-deck`."""
import io

from pptx import Presentation

from sonaloop import _deck, _pptx


def _slide_text(slide) -> str:
    return " ".join(r.text for sh in slide.shapes if sh.has_text_frame
                    for p in sh.text_frame.paragraphs for r in p.runs)


def test_sample_slides_cover_every_layout():
    """The placeholder deck demonstrates the FULL taxonomy — one sample per layout, in deck order."""
    assert [s["kind"] for s in _deck.SAMPLE_SLIDES] == [l["key"] for l in _deck.LAYOUTS]
    assert len(_deck.LAYOUTS) >= 20  # cover … closing + brand/charts/table layouts + fallbacks


def test_master_template_renders_one_slide_per_sample():
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == len(_deck.SAMPLE_SLIDES)


def test_every_layout_paints_visible_content():
    """Each layout leaves real text on its slide (the image fallback paints a placeholder panel)."""
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    for spec, slide in zip(_deck.SAMPLE_SLIDES, prs.slides):
        text = _slide_text(slide)
        assert text.strip(), f"layout {spec['kind']!r} rendered an empty slide"
        assert len(slide.shapes) >= 3, f"layout {spec['kind']!r} painted too few shapes"


def test_key_sample_copy_lands_on_the_right_slides():
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    by_kind = {spec["kind"]: _slide_text(slide)
               for spec, slide in zip(_deck.SAMPLE_SLIDES, prs.slides)}
    assert "Healthy eating without effort" in by_kind["cover"]
    assert "RESEARCH REPORT" in by_kind["cover"]            # eyebrow is uppercased
    assert "INSIGHT" in by_kind["insight"]                  # tone chip
    assert "RECOMMENDATION" in by_kind["recommendation"]
    assert "RISK" in by_kind["risk"]
    assert "Mehmet" in by_kind["quote"]
    assert "SUPPORT" in by_kind["voices"]                   # sentiment chip
    assert "124" in by_kind["stats"]


def test_brand_assets_land_as_pictures():
    """Cover (logo + canvas band), canvas-section (full-bleed art), pillars (icon chips) and
    closing (logo) embed the vendored _deck_assets rasters as real PICTURE shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    pics = {spec["kind"]: sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
            for spec, slide in zip(_deck.SAMPLE_SLIDES, prs.slides)}
    assert pics["cover"] == 2          # logo mark + canvas band
    assert pics["canvas-section"] == 1  # full-bleed canvas
    n_pillars = len(next(s for s in _deck.SAMPLE_SLIDES if s["kind"] == "pillars")["items"])
    assert pics["pillars"] == n_pillars
    assert pics["closing"] == 1        # logo mark


def test_brand_slides_carry_their_copy():
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    by_kind = {spec["kind"]: _slide_text(slide)
               for spec, slide in zip(_deck.SAMPLE_SLIDES, prs.slides)}
    assert "sona" in by_kind["cover"] and "loop" in by_kind["cover"]  # wordmark
    assert "Three product directions" in by_kind["canvas-section"]
    assert "Grounded personas" in by_kind["pillars"]


def test_table_slide_is_a_native_table():
    """The table layout lands as a real GraphicFrame table — header + one row per data row."""
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    spec, slide = next((s, sl) for s, sl in zip(_deck.SAMPLE_SLIDES, prs.slides)
                       if s["kind"] == "table")
    tbl = next(sh.table for sh in slide.shapes if sh.has_table)
    assert len(tbl.columns) == len(spec["columns"])
    assert len(tbl.rows) == len(spec["rows"]) + 1
    assert tbl.cell(0, 0).text == spec["columns"][0]
    assert tbl.cell(1, 0).text == spec["rows"][0][0]


def test_charts_slide_carries_two_native_charts():
    """The charts layout paints both chart slots (shapes, not a screenshot)."""
    data = _pptx.render(_deck.SAMPLE_SLIDES, title=_deck.DECK_TITLE)
    prs = Presentation(io.BytesIO(data))
    spec, slide = next((s, sl) for s, sl in zip(_deck.SAMPLE_SLIDES, prs.slides)
                       if s["kind"] == "charts")
    text = " ".join(r.text for sh in slide.shapes if sh.has_text_frame
                    for p in sh.text_frame.paragraphs for r in p.runs)
    for it in spec["items"]:
        assert it["title"] in text
    # both slots draw real shape clusters under their captions
    assert len(slide.shapes) > 20


def test_unknown_brand_assets_degrade_silently():
    """A canvas/icon name outside the curated set must not crash the renderer."""
    slides = [{"kind": "cover", "title": "T", "logo": True, "canvas": "nope"},
              {"kind": "canvas-section", "title": "S", "canvas": "nope"},
              {"kind": "pillars", "heading": "H", "items": [{"icon": "nope", "title": "x"}]}]
    prs = Presentation(io.BytesIO(_pptx.render(slides)))
    assert len(prs.slides) == 3


def test_palette_and_tones_come_from_the_design_system():
    """The renderer's brand constants are the vendored deck palette — no drift-prone hex literals."""
    assert _pptx._ACCENT == _deck.PALETTE["accent"]
    assert _pptx._SERIES == _deck.PALETTE["series"]
    assert set(_deck.TONES) == {"insight", "recommendation", "risk"}


def test_template_deck_cli_writes_the_demo_deck(tmp_path, capsys):
    from sonaloop.cli import main
    out = tmp_path / "demo.pptx"
    main(["template-deck", "--out", str(out)])
    assert out.exists()
    prs = Presentation(io.BytesIO(out.read_bytes()))
    assert len(prs.slides) == len(_deck.SAMPLE_SLIDES)
