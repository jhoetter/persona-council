"""Charts are design-system components (sonaloop._charts, vendored from sonaloop-design) embedded in
reports via the `chart` figure-kind: bar · pie · effort·impact. (Phase 2 — spec/unified-synthesis-report.)"""
from __future__ import annotations

import base64

from sonaloop import _charts
from sonaloop.web._report import render_report


def test_chart_components_render_and_are_empty_safe():
    bar = _charts.bar_chart([{"label": "Plan", "value": 8}, {"label": "Cook", "value": 3}])
    assert "sl-bars" in bar and "sl-bar__fill" in bar and "--v:100.0%" in bar  # max value → full bar
    pie = _charts.pie_chart([{"label": "Support", "value": 12}, {"label": "Oppose", "value": 4}])
    assert "conic-gradient(" in pie and "sl-pie--donut" in pie and "75%" in pie  # 12/16
    ei = _charts.effort_impact([{"label": "Auto list", "x": 2, "y": 5}], x_label="Effort", y_label="Value")
    assert "sl-quad__dot" in ei and "sl-legend__num" in ei and "Quick wins" in ei
    stacked = _charts.stacked_bar_chart([
        {"label": "Pricing", "segments": [{"label": "For", "value": 6}, {"label": "Against", "value": 2}]},
        {"label": "Support", "segments": [{"label": "For", "value": 2}, {"label": "Against", "value": 5}]}])
    # series colour keyed by segment label → one shared legend with 2 series (not 4)
    assert stacked.count("sl-bar__fill--stack") == 2 and "sl-bar__seg" in stacked
    assert "sl-legend--row" in stacked and stacked.count("sl-legend__sw") == 2
    gauge = _charts.gauge_chart([{"label": "Confidence", "value": 72}, {"label": "Done", "value": 9, "max": 12}])
    assert "sl-gauge" in gauge and "--p:72.0" in gauge and "72%" in gauge and "75%" in gauge  # 9/12
    assert "9 / 12" in gauge  # non-100 max shows the raw ratio
    diverging = _charts.diverging_bar_chart(
        [{"label": "Pricing", "positive": 6, "negative": 2}], positive_label="For", negative_label="Against")
    assert "sl-dbar__pos" in diverging and "sl-dbar__neg" in diverging and "+6 · −2" in diverging
    heat = _charts.heatmap_chart(["Cost", "Reach"], [{"label": "Plan A", "values": [2, 5]}])
    assert "sl-heat" in heat and "color-mix(in srgb" in heat and "grid-template-columns" in heat
    dots = _charts.dot_plot_chart([{"label": "Trust", "values": [2, 3, 3, 4, 5]}])
    assert dots.count("sl-dot-pt") == 5 and "sl-dot-mean" in dots and "3.4" in dots  # mean of values
    line = _charts.line_chart([{"label": "Conf", "points": [2, 3, 5, 4, 6]}], labels=["R1", "R2", "R3", "R4", "R5"])
    assert "<svg" in line and "sl-line__path" in line and "<polyline" in line and "sl-line__labels" in line
    # empty / unscored inputs never fabricate a chart
    assert _charts.bar_chart([]) == "" and _charts.pie_chart([{"label": "x", "value": 0}]) == ""
    assert _charts.effort_impact([{"label": "no score", "x": None, "y": 2}]) == ""
    assert _charts.stacked_bar_chart([]) == "" and _charts.gauge_chart([]) == ""
    assert _charts.diverging_bar_chart([]) == "" and _charts.heatmap_chart([], []) == ""
    assert _charts.dot_plot_chart([]) == "" and _charts.line_chart([{"label": "x", "points": [1]}]) == ""


def test_new_chart_components_render_and_are_empty_safe():
    """The Linear wave: burnup · stacked_area · column · strip · progress_strip · stats + micro
    indicators (sparkline / progress_pie)."""
    burnup = _charts.burnup_chart([{"label": "Done", "points": [0, 2, 5, 8]}],
                                  labels=["W1", "W2", "W3", "W4"], target=12, now=2)
    assert "sl-line__area" in burnup and "sl-line__ref" in burnup  # band fill + dotted ideal
    assert "sl-line__now" in burnup and "sl-burnup__future" in burnup and "sl-burnup__hatch" in burnup
    area = _charts.stacked_area_chart([{"label": "For", "points": [2, 4, 6]},
                                       {"label": "Against", "points": [3, 2, 1]}])
    assert area.count("sl-area__band") == 2 and area.count("sl-area__edge") == 2
    assert area.count("sl-legend__item") == 2  # bands always get a legend
    col = _charts.column_chart([{"label": "1", "value": 2}, {"label": "2", "value": 5}], table=True)
    assert "sl-cols-axis" in col and col.count('sl-col"') == 2 and "sl-chart__table" in col
    colseg = _charts.column_chart([{"label": "Todo", "segments": [
        {"label": "High", "value": 7}, {"label": "Low", "value": 4}]}])
    assert "sl-col__bar--stack" in colseg and colseg.count("sl-legend__item") == 2
    strip = _charts.strip_chart([{"label": "WTP", "values": [9, 15, 29]}], unit="€")
    assert strip.count("sl-dot-pt") == 3 and "9€" in strip and "29€" in strip and "19€" in strip  # min/mid/max
    pstrip = _charts.progress_strip([{"label": "Validated", "value": 9}, {"label": "Open", "value": 3}])
    assert pstrip.count("sl-pstrip__seg") == 2 and "· 75%" in pstrip
    stats = _charts.stats_chart([{"label": "Personas", "value": 16},
                                 {"label": "Agreement", "value": "72%", "sub": "+9 vs R1", "color": "var(--c1)"}])
    assert stats.count('sl-kpi"') == 2 and "72%" in stats and "sl-kpi__sub" in stats and "sl-kpi__sw" in stats
    spark = _charts.sparkline([3, 5, 4, 6])
    assert "sl-spark__line" in spark and "sl-spark__fill" in spark
    mpie = _charts.progress_pie(11, 16)
    assert "sl-mpie" in mpie and "--p:68.8" in mpie
    # empty / unscored inputs never fabricate a chart
    assert _charts.burnup_chart([]) == "" and _charts.stacked_area_chart([{"label": "x", "points": [1]}]) == ""
    assert _charts.column_chart([]) == "" and _charts.strip_chart([{"label": "x", "values": []}]) == ""
    assert _charts.progress_strip([{"label": "x", "value": 0}]) == "" and _charts.stats_chart([]) == ""
    assert _charts.sparkline([1]) == ""


def test_chart_escapes_labels():
    out = _charts.bar_chart([{"label": "<script>", "value": 1}])
    assert "<script>" not in out and "&lt;script&gt;" in out


def _report_with_figures(figures):
    return {"scope": "project", "title": "R — Report", "created_at": "2026-06-08T00:00:00+00:00",
            "lead": "", "graph_snapshot": None,
            "sections": [{"id": "s1", "heading": "Findings", "markdown": "Body.",
                          "citations": [], "figures": figures, "source_study_ids": []}]}


def test_report_embeds_author_supplied_pie_and_bar(store):
    html = render_report(_report_with_figures([
        {"kind": "chart", "of": "pie", "series": [{"label": "A", "value": 3}, {"label": "B", "value": 1}], "caption": "Split"},
        {"kind": "chart", "of": "bar", "series": [{"label": "X", "value": 5}]},
    ]), store)
    assert "sl-pie" in html and "conic-gradient" in html and "sl-bars" in html


def test_report_embeds_author_supplied_stacked_bar_and_gauge(store):
    html = render_report(_report_with_figures([
        {"kind": "chart", "of": "stacked_bar", "caption": "Stance per theme", "series": [
            {"label": "Pricing", "segments": [{"label": "For", "value": 6}, {"label": "Against", "value": 2}]}]},
        {"kind": "chart", "of": "gauge", "series": [{"label": "Confidence", "value": 72}]},
    ]), store)
    assert "sl-bar__fill--stack" in html and "sl-bar__seg" in html
    assert "sl-gauge" in html and "72%" in html


def test_section_validator_preserves_gauge_max_and_stacked_segments():
    from sonaloop.llm_simulation._validators import validate_synthesis_section_payload
    out = validate_synthesis_section_payload({"markdown": "Body.", "figures": [
        {"kind": "chart", "of": "gauge", "series": [{"label": "Confidence", "value": 72, "max": 100}]},
        {"kind": "chart", "of": "stacked_bar", "series": [
            {"label": "Pricing", "segments": [{"label": "For", "value": 6}, {"label": "Against", "value": 2}]}]},
    ]})
    gauge, stacked = out["figures"]
    assert gauge["series"][0]["max"] == 100  # gauge scale survives sanitization
    segs = stacked["series"][0]["segments"]
    assert [s["label"] for s in segs] == ["For", "Against"] and segs[0]["value"] == 6


def test_report_embeds_diverging_heatmap_dotplot_line(store):
    html = render_report(_report_with_figures([
        {"kind": "chart", "of": "diverging_bar", "series": [{"label": "Pricing", "positive": 6, "negative": 2}]},
        {"kind": "chart", "of": "heatmap", "columns": ["Cost", "Reach"],
         "series": [{"label": "Plan A", "values": [2, 5]}]},
        {"kind": "chart", "of": "dot_plot", "series": [{"label": "Trust", "values": [2, 3, 4, 5]}]},
        {"kind": "chart", "of": "line", "labels": ["R1", "R2", "R3"],
         "series": [{"label": "Conf", "points": [2, 4, 6]}]},
    ]), store)
    assert "sl-dbar__pos" in html  # diverging
    assert "sl-heat" in html and "color-mix(in srgb" in html  # heatmap
    assert "sl-dot-track" in html  # dot plot
    assert "sl-line__path" in html and "<svg" in html  # line


def test_section_validator_preserves_new_chart_fields():
    from sonaloop.llm_simulation._validators import validate_synthesis_section_payload
    out = validate_synthesis_section_payload({"markdown": "Body.", "figures": [
        {"kind": "chart", "of": "diverging_bar", "series": [{"label": "Pricing", "positive": 6, "negative": 2}]},
        {"kind": "chart", "of": "heatmap", "columns": ["Cost", "Reach"],
         "series": [{"label": "Plan A", "values": [2, 5]}]},
        {"kind": "chart", "of": "dot_plot", "series": [{"label": "Trust", "values": [2, 3, 4, 5]}]},
        {"kind": "chart", "of": "line", "labels": ["R1", "R2"], "series": [{"label": "Conf", "points": [2, 4]}]},
    ]})
    diverging, heat, dots, line = out["figures"]
    assert diverging["series"][0]["positive"] == 6 and diverging["series"][0]["negative"] == 2
    assert heat["columns"] == ["Cost", "Reach"] and heat["series"][0]["values"] == [2, 5]
    assert dots["series"][0]["values"] == [2, 3, 4, 5]
    assert line["labels"] == ["R1", "R2"] and line["series"][0]["points"] == [2, 4]


def test_report_embeds_new_chart_kinds(store):
    html = render_report(_report_with_figures([
        {"kind": "chart", "of": "burnup", "labels": ["W1", "W2", "W3"], "target": 12, "now": 1,
         "series": [{"label": "Done", "points": [0, 3, 7]}]},
        {"kind": "chart", "of": "stacked_area", "series": [
            {"label": "For", "points": [2, 4, 6]}, {"label": "Against", "points": [3, 2, 1]}]},
        {"kind": "chart", "of": "column", "table": True, "series": [{"label": "1", "value": 2}, {"label": "2", "value": 5}]},
        {"kind": "chart", "of": "strip", "unit": "€", "series": [{"label": "WTP", "values": [9, 15, 29]}]},
        {"kind": "chart", "of": "progress_strip", "series": [{"label": "Open", "value": 3}, {"label": "Done", "value": 9}]},
        {"kind": "chart", "of": "stats", "series": [{"label": "Personas", "value": 16, "sub": "of 20"}]},
    ]), store)
    assert "sl-line__ref" in html and "sl-burnup__future" in html  # burnup target + future region
    assert "sl-area__band" in html                                  # stacked area
    assert "sl-cols" in html and "sl-chart__table" in html          # column + its table
    assert "sl-dot-pt" in html and "29€" in html                    # strip with unit
    assert "sl-pstrip__seg" in html                                 # progress strip
    assert "sl-kpi__val" in html and "sl-kpi__sub" in html          # stats tiles


def test_section_validator_preserves_linear_wave_fields():
    from sonaloop.llm_simulation._validators import validate_synthesis_section_payload
    out = validate_synthesis_section_payload({"markdown": "Body.", "figures": [
        {"kind": "chart", "of": "burnup", "labels": ["W1", "W2"], "target": 12, "now": 1,
         "series": [{"label": "Done", "points": [0, 3]}]},
        {"kind": "chart", "of": "strip", "min": 0, "max": 50, "unit": "€",
         "series": [{"label": "WTP", "values": [9, 15]}]},
        {"kind": "chart", "of": "column", "table": True, "series": [{"label": "1", "value": 2}]},
        {"kind": "chart", "of": "stats", "series": [{"label": "Personas", "value": 16, "sub": "of 20"}]},
    ]})
    burnup, strip, column, stats = out["figures"]
    assert burnup["target"] == 12 and burnup["now"] == 1 and burnup["labels"] == ["W1", "W2"]
    assert strip["min"] == 0 and strip["max"] == 50 and strip["unit"] == "€"
    assert column["table"] is True
    assert stats["series"][0]["sub"] == "of 20"
    # booleans never sneak in as numbers; junk unit/table values are dropped or normalized
    out2 = validate_synthesis_section_payload({"markdown": "B.", "figures": [
        {"kind": "chart", "of": "burnup", "target": True, "series": [{"label": "D", "points": [0, 3]}]}]})
    assert "target" not in out2["figures"][0]


def test_chart_catalogue_is_the_single_source_of_truth():
    """The agent-facing catalogue, the renderer registry, and the report dispatch never drift:
    every author-renderable `of` is documented, and every documented `of` actually renders."""
    from sonaloop import services
    from sonaloop.charts_catalogue import _RENDER, render_chart
    cat = services.suggest_chart_kinds()
    kinds = {k["of"] for k in cat["items"]}
    # catalogue = the renderable kinds + the source-driven effort_impact, and nothing else
    assert kinds == set(_RENDER) | {"effort_impact"}
    # each documented entry carries the agent's "when to use" guidance
    assert all(k.get("when_to_use") for k in cat["items"])
    # every renderable kind produces a chart from a minimal series; unknown/source-driven → ""
    sample = {
        "bar": [{"label": "A", "value": 3}], "pie": [{"label": "A", "value": 3}],
        "stacked_bar": [{"label": "A", "segments": [{"label": "x", "value": 2}]}],
        "diverging_bar": [{"label": "A", "positive": 2, "negative": 1}],
        "gauge": [{"label": "A", "value": 50}], "dot_plot": [{"label": "A", "values": [1, 2, 3]}],
        "heatmap": [{"label": "A", "values": [2]}], "line": [{"label": "A", "points": [1, 2, 3]}],
        "burnup": [{"label": "A", "points": [1, 2, 3]}],
        "stacked_area": [{"label": "A", "points": [1, 2, 3]}],
        "column": [{"label": "A", "value": 3}],
        "strip": [{"label": "A", "values": [1, 2, 3]}],
        "progress_strip": [{"label": "A", "value": 3}],
        "stats": [{"label": "A", "value": 3}],
    }
    for of in _RENDER:
        fig = {"columns": ["c"]} if of == "heatmap" else {"labels": ["a", "b", "c"]} if of in ("line", "burnup", "stacked_area") else {}
        assert render_chart(of, fig, sample[of]).startswith('<figure class="sl-chart">'), of
    assert render_chart("effort_impact", {}, []) == "" and render_chart("bogus", {}, []) == ""


def test_report_embeds_synthesis_effort_impact_2x2(store):
    from sonaloop import services
    syn = services.record_synthesis(
        "Solutions", "hmw", [], {"findings": [
            {"kind": "recommendation", "text": "Auto shopping list", "score": {"effort": 2, "value": 5}},
            {"kind": "recommendation", "text": "Hire a coach", "score": {"effort": 4, "value": 2}}]},
        store=store)
    html = render_report(_report_with_figures(
        [{"kind": "chart", "of": "effort_impact", "source_id": syn["id"], "caption": "Leverage"}]), store)
    assert "sl-quad" in html and "sl-quad__dot" in html and "sl-legend__num" in html


def _has_chart(shape):
    return bool(getattr(shape, "has_chart", False))


# a minimal valid 1×1 PNG (so add_picture accepts it)
_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg==")


def test_pptx_embeds_image_figure(tmp_path):
    """Prototype screenshots / image assets must travel in the deck (the PDF loads them by URL; a PPTX
    has to carry the bytes). An image slide embeds the picture."""
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from sonaloop import _pptx
    png = tmp_path / "shot.png"; png.write_bytes(_PNG_1X1)
    data = _pptx.render([{"kind": "image", "heading": "Prototype", "num": "03",
                          "image": str(png), "caption": "v0.2"}], title="R")
    prs = Presentation(io.BytesIO(data))
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sl in prs.slides for sh in sl.shapes)


def test_write_export_bytes_round_trip(tmp_path):
    """Binary exports (pdf/pptx) write to disk — the MCP export tool / CLI return this path."""
    from sonaloop import services
    out = services.write_export_bytes(b"%PDF-x", tmp_path / "r.pdf")
    assert out.endswith("r.pdf")
    with open(out, "rb") as fh:
        assert fh.read() == b"%PDF-x"


def test_report_section_image_figure_becomes_image_slide(store, tmp_path, monkeypatch):
    """A section's prototype/asset figure resolves to a file and exports as an image slide."""
    import io
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from sonaloop import services, assets
    monkeypatch.setattr(assets, "ASSETS_DIR", tmp_path / "assets")
    monkeypatch.setattr("sonaloop.config.DATA_DIR", tmp_path)
    aid = assets.put_asset(_PNG_1X1, "png")
    rep = {"id": "rimg", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-08T00:00:00+00:00", "lead": "", "council_ids": [],
           "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": "s1", "heading": "Deliver", "markdown": "Body.", "citations": [],
                         "source_study_ids": [], "figures": [{"kind": "asset", "id": aid, "caption": "Shot"}]}]}
    store.upsert_synthesis(rep)
    data = services.export_synthesis_pptx("rimg", store=store)
    prs = Presentation(io.BytesIO(data))
    assert any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sl in prs.slides for sh in sl.shapes)


def test_pptx_export_convergence_synthesis(store):
    """A structured synthesis exports to a native .pptx deck: title + analytic-layer slides + the
    effort·impact 2×2 as a native scatter chart."""
    import io
    from pptx import Presentation
    from sonaloop import services
    syn = services.record_synthesis(
        "Solutions", "hmw", [], {
            "gesamtbild": "G" * 80, "positionierung": "P" * 80,
            "findings": [
                {"kind": "key_problem", "text": "Evenings are the bottleneck"},
                {"kind": "recommendation", "text": "Auto shopping list", "score": {"effort": 2, "value": 5}},
                {"kind": "recommendation", "text": "Hire a coach", "score": {"effort": 4, "value": 2}},
                {"kind": "open_question", "text": "Does it survive week 3?"}]},
        store=store)
    data = services.export_synthesis_pptx(syn["id"], store=store)
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 5  # title + big picture + positioning + key problems + recommendations + open Qs
    # the effort·impact 2×2 is drawn from native shapes (frame, dashed cross, numbered dots, legend)
    assert sum(len(sl.shapes) for sl in prs.slides) > 20


def test_pptx_export_project_report_with_chart(store):
    """A project report exports one slide per section + author-supplied charts as native pptx charts."""
    import io
    from pptx import Presentation
    from sonaloop import services
    rep = {"id": "rep1", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-08T00:00:00+00:00", "lead": "How the understanding was built.",
           "council_ids": [], "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": "s1", "heading": "Findings", "markdown": "- Point one\n- Point two",
                         "citations": [], "source_study_ids": [],
                         "figures": [{"kind": "chart", "of": "pie",
                                      "series": [{"label": "A", "value": 3}, {"label": "B", "value": 1}]}]}]}
    store.upsert_synthesis(rep)
    data = services.export_synthesis_pptx("rep1", store=store)
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= 2  # title + section
    assert any(_has_chart(sh) for sl in prs.slides for sh in sl.shapes)  # the pie


def test_figure_to_chart_maps_every_design_system_kind(store):
    """The PPTX bridge mirrors the design-system chart of-kinds (charts_catalogue), so a deck shows
    the same chart the report/PDF does — one neutral _pptx model per author-supplied `of`."""
    from sonaloop.services._synthesis_pptx import _figure_to_chart
    from sonaloop.charts_catalogue import _RENDER
    figs = {
        "bar": {"of": "bar", "series": [{"label": "A", "value": 3}]},
        "pie": {"of": "pie", "series": [{"label": "A", "value": 3}]},
        "stacked_bar": {"of": "stacked_bar", "series": [{"label": "A", "segments": [{"label": "x", "value": 2}]}]},
        "diverging_bar": {"of": "diverging_bar", "series": [{"label": "A", "positive": 2, "negative": 1}]},
        "gauge": {"of": "gauge", "series": [{"label": "A", "value": 50, "max": 100}]},
        "dot_plot": {"of": "dot_plot", "series": [{"label": "A", "values": [1, 2, 3]}]},
        "heatmap": {"of": "heatmap", "columns": ["c"], "series": [{"label": "A", "values": [2]}]},
        "line": {"of": "line", "labels": ["a", "b", "c"], "series": [{"label": "A", "points": [1, 2, 3]}]},
        "burnup": {"of": "burnup", "target": 5, "series": [{"label": "A", "points": [1, 2, 3]}]},
        "stacked_area": {"of": "stacked_area", "series": [{"label": "A", "points": [1, 2, 3]}]},
        "column": {"of": "column", "series": [{"label": "A", "value": 3}]},
        "strip": {"of": "strip", "unit": "€", "series": [{"label": "A", "values": [1, 2, 3]}]},
        "progress_strip": {"of": "progress_strip", "series": [{"label": "A", "value": 3}]},
        "stats": {"of": "stats", "series": [{"label": "A", "value": 3, "sub": "of 5"}]},
    }
    # every author-renderable design-system kind has a PPTX bridge mapping (parity with the HTML side)
    assert set(figs) == set(_RENDER)
    # burnup reuses the native line model (dashed Target series); strip reuses dot_plot (real scale)
    expect = {**{of: of for of in figs}, "burnup": "line", "strip": "dot_plot"}
    for of, fig in figs.items():
        model = _figure_to_chart({"kind": "chart", **fig}, store)
        assert model and model["type"] == expect[of], of
    assert _figure_to_chart({"kind": "chart", "of": "burnup", "series": [{"label": "A", "points": [1, 2]}]},
                            store)["target"] is None  # target optional
    assert _figure_to_chart({"kind": "chart", "of": "strip", "series": [{"label": "A", "values": [3, 9]}]},
                            store)["min"] == 3  # scale derived from the data when not fixed
    assert _figure_to_chart({"kind": "asset", "id": "x"}, store) is None  # non-chart
    assert _figure_to_chart({"kind": "chart", "of": "line",  # a 1-point "line" can't be drawn
                             "series": [{"label": "A", "points": [1]}]}, store) is None


def test_pptx_export_renders_all_new_chart_kinds(store):
    """Every new chart kind exports into the deck — native charts (gauge/line) and shape-drawn charts
    (stacked/diverging/dot/heatmap) — without error and as real shapes."""
    import io
    from pptx import Presentation
    from sonaloop import services
    figures = [
        {"kind": "chart", "of": "stacked_bar", "series": [
            {"label": "Pricing", "segments": [{"label": "For", "value": 6}, {"label": "Against", "value": 2}]}]},
        {"kind": "chart", "of": "diverging_bar", "series": [{"label": "Pricing", "positive": 6, "negative": 2}]},
        {"kind": "chart", "of": "gauge", "series": [{"label": "Confidence", "value": 72}]},
        {"kind": "chart", "of": "dot_plot", "series": [{"label": "Trust", "values": [2, 3, 4, 5]}]},
        {"kind": "chart", "of": "heatmap", "columns": ["Cost", "Reach"],
         "series": [{"label": "Plan A", "values": [2, 5]}]},
        {"kind": "chart", "of": "line", "labels": ["R1", "R2", "R3"], "series": [{"label": "Conf", "points": [2, 4, 6]}]},
    ]
    rep = {"id": "repall", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-08T00:00:00+00:00", "lead": "", "council_ids": [],
           "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": f"s{i}", "heading": f"Section {i}", "markdown": "Body.", "citations": [],
                         "source_study_ids": [], "figures": [fig]} for i, fig in enumerate(figures)]}
    store.upsert_synthesis(rep)
    data = services.export_synthesis_pptx("repall", store=store)
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= len(figures) + 1  # title + one slide per section
    assert sum(_has_chart(sh) for sl in prs.slides for sh in sl.shapes) >= 2  # native gauge + line
    assert sum(len(sl.shapes) for sl in prs.slides) > 30  # the shape-drawn charts add many shapes


def test_pptx_export_renders_linear_wave_chart_kinds(store):
    """The Linear wave exports too: burnup (line + dashed Target), stacked_area (native AREA_STACKED),
    column / progress_strip / stats (shape-drawn), strip (dot_plot scale)."""
    import io
    from pptx import Presentation
    from sonaloop import services
    figures = [
        {"kind": "chart", "of": "burnup", "labels": ["W1", "W2", "W3"], "target": 12,
         "series": [{"label": "Done", "points": [0, 3, 7]}]},
        {"kind": "chart", "of": "stacked_area", "series": [
            {"label": "For", "points": [2, 4, 6]}, {"label": "Against", "points": [3, 2, 1]}]},
        {"kind": "chart", "of": "column", "series": [
            {"label": "Todo", "segments": [{"label": "High", "value": 7}, {"label": "Low", "value": 4}]},
            {"label": "Done", "segments": [{"label": "High", "value": 5}, {"label": "Low", "value": 2}]}]},
        {"kind": "chart", "of": "strip", "unit": "€", "series": [{"label": "WTP", "values": [9, 15, 29]}]},
        {"kind": "chart", "of": "progress_strip", "series": [
            {"label": "Open", "value": 3}, {"label": "Done", "value": 9}]},
        {"kind": "chart", "of": "stats", "series": [
            {"label": "Personas", "value": 16}, {"label": "Agreement", "value": "72%", "sub": "+9"}]},
    ]
    rep = {"id": "replw", "title": "Demo — Report", "scope": "project", "project_id": "",
           "created_at": "2026-06-10T00:00:00+00:00", "lead": "", "council_ids": [],
           "findings": [], "statements": [], "prompts": [], "graph_snapshot": None,
           "sections": [{"id": f"s{i}", "heading": f"Section {i}", "markdown": "Body.", "citations": [],
                         "source_study_ids": [], "figures": [fig]} for i, fig in enumerate(figures)]}
    store.upsert_synthesis(rep)
    data = services.export_synthesis_pptx("replw", store=store)
    assert data[:2] == b"PK"
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) >= len(figures) + 1
    assert sum(_has_chart(sh) for sl in prs.slides for sh in sl.shapes) >= 2  # native burnup-line + area
    texts = " ".join(sh.text_frame.text for sl in prs.slides for sh in sl.shapes if sh.has_text_frame)
    assert "72%" in texts and "Personas" in texts  # the stats tiles made it into the deck


def test_line_chart_target_reference_reaches_render(store):
    """The aesthetics pass: a plain line figure can carry target= for a dotted reference line."""
    html = render_report(_report_with_figures([
        {"kind": "chart", "of": "line", "target": 7, "labels": ["R1", "R2", "R3"],
         "series": [{"label": "Conf", "points": [2, 4, 6]}]},
    ]), store)
    assert "sl-line__ref" in html and "sl-line__grid" in html and "sl-line__area" in html
