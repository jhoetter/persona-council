"""First-slide PPTX preview rasterizer (ux-contract §10 W6 — "schon die Titel-Slide anzeigen").

The productized V11 verification renderer: PIL drawing over python-pptx shape geometry —
solid fills, rounded rects/ovals, pictures, tables, connector lines and word-wrapped text
runs. Good enough to SHOW the title slide on a document file card; deliberately NOT a print
path (no gradients/effects/embedded-font fidelity — PowerPoint owns that).

The one entry point is `render_first_slide(data) -> bytes | None`: PNG bytes scaled to
PREVIEW_WIDTH, or None when the bytes aren't a renderable deck (corrupt file, python-pptx
missing) — callers degrade to the extension badge, never to an error. The asset seam
(services._project_assets._write_asset_preview) writes the result as
`<sha>.preview.png` beside the binary in the content-addressed store.
"""
from __future__ import annotations

import io

PREVIEW_WIDTH = 640          # the card stage consumes ~half; 2x for crisp retina thumbs
_EMU_PER_PX = 914400 / 96    # render geometry at 96 dpi, then scale down

# Font candidates per face — first hit wins; PIL's built-in bitmap face is the last resort
# (Linux CI/laptops ship DejaVu; mac/Windows fall through to the system paths).
_FONT_PATHS = {
    "regular": ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/System/Library/Fonts/Helvetica.ttc", "C:/Windows/Fonts/arial.ttf"),
    "bold": ("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
             "/System/Library/Fonts/Helvetica.ttc", "C:/Windows/Fonts/arialbd.ttf"),
    "mono": ("/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
             "/System/Library/Fonts/Menlo.ttc", "C:/Windows/Fonts/consola.ttf"),
}
_fonts: dict = {}


def _font(size_pt: float, bold: bool = False, mono: bool = False):
    from PIL import ImageFont
    px = max(6, int(size_pt * 96 / 72))
    key = (px, bold, mono)
    if key not in _fonts:
        face = "mono" if mono else ("bold" if bold else "regular")
        font = None
        for path in _FONT_PATHS[face]:
            try:
                font = ImageFont.truetype(path, px)
                break
            except OSError:
                continue
        if font is None:
            try:
                font = ImageFont.load_default(size=px)
            except TypeError:          # older Pillow: unsized bitmap fallback
                font = ImageFont.load_default()
        _fonts[key] = font
    return _fonts[key]


def _px(emu) -> int:
    return int(round(emu / _EMU_PER_PX))


def _solid_color(fill, default=None):
    try:
        if fill.type is not None and str(fill.type) != "MSO_FILL_TYPE.BACKGROUND":
            return "#" + str(fill.fore_color.rgb)
    except Exception:
        pass
    return default


def _run_color(r, default="#1A1815"):
    try:
        if r.font.color and r.font.color.rgb is not None:
            return "#" + str(r.font.color.rgb)
    except Exception:
        pass
    return default


def _wrap_runs(draw, runs, width):
    """Greedy word wrap over styled runs [(text, font, color)] -> list of line run-lists."""
    lines, cur, cur_w = [], [], 0.0
    for text, f, c in runs:
        for piece in text.split(" "):
            if not piece:
                continue
            w = draw.textlength(piece, font=f)
            add = w + (draw.textlength(" ", font=f) if cur else 0)
            if cur and cur_w + add > width:
                lines.append(cur)
                cur, cur_w = [(piece, f, c)], w
            else:
                cur.append((piece, f, c))
                cur_w += add
    if cur:
        lines.append(cur)
    return lines


def _draw_text_frame(draw, tf, box):
    x0, y0, w, hgt = box
    pad = 4
    blocks = []
    for p in tf.paragraphs:
        runs = []
        for r in p.runs:
            size = r.font.size.pt if r.font.size else 14
            mono = "mono" in (r.font.name or "").lower()
            f = _font(size, bold=bool(r.font.bold), mono=mono)
            runs.append((r.text, f, _run_color(r)))
        if not runs:
            blocks.append((p, [], 6.0, 0.0, 0.0))
            continue
        lines = _wrap_runs(draw, runs, max(20, w - 2 * pad))
        lh = max((f.size for _, f, _ in runs), default=12) * 1.25
        sb = (p.space_before.pt * 96 / 72) if p.space_before else 0
        sa = (p.space_after.pt * 96 / 72) if p.space_after else 0
        blocks.append((p, lines, lh, sb, sa))
    total = sum(sb + len(lines) * lh + sa for _, lines, lh, sb, sa in blocks)
    anchor = None
    try:
        anchor = tf.vertical_anchor
    except Exception:
        pass
    y = y0 + ((hgt - total) / 2 if anchor is not None and str(anchor).startswith("MIDDLE")
              and total < hgt else pad)
    for p, lines, lh, sb, sa in blocks:
        if not lines:
            y += lh
            continue
        y += sb
        for line in lines:
            lw = sum(draw.textlength(tk + " ", font=f) for tk, f, _ in line)
            lw -= draw.textlength(" ", font=line[-1][1]) if line else 0
            align = str(p.alignment) if p.alignment else ""
            if "CENTER" in align:
                x = x0 + (w - lw) / 2
            elif "RIGHT" in align:
                x = x0 + w - lw - pad
            else:
                x = x0 + pad
            for tk, f, c in line:
                draw.text((x, y), tk, font=f, fill=c)
                x += draw.textlength(tk + " ", font=f)
            y += lh
        y += sa


def _line_color(sh, default=None):
    try:
        if sh.line.color and sh.line.color.rgb is not None:
            return "#" + str(sh.line.color.rgb)
    except Exception:
        pass
    return default


def _render_shape(img, draw, sh):
    from PIL import Image
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    try:
        x, y, w, hgt = _px(sh.left), _px(sh.top), _px(sh.width), _px(sh.height)
    except Exception:
        return
    st = sh.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            im = Image.open(io.BytesIO(sh.image.blob)).convert("RGBA")
            iw, ih = im.size
            im = im.crop((int(iw * sh.crop_left), int(ih * sh.crop_top),
                          int(iw * (1 - sh.crop_right)), int(ih * (1 - sh.crop_bottom))))
            im = im.resize((max(1, w), max(1, hgt)))
            img.paste(im, (x, y), im)
        except Exception:
            draw.rectangle([x, y, x + w, y + hgt], outline="#cccccc")
        return
    if st == MSO_SHAPE_TYPE.TABLE:
        tbl = sh.table
        col_w = [_px(c.width) for c in tbl.columns]
        row_h = [_px(r.height) for r in tbl.rows]
        ty = y
        for ri, row in enumerate(tbl.rows):
            tx = x
            for ci, cell in enumerate(row.cells):
                draw.rectangle([tx, ty, tx + col_w[ci], ty + row_h[ri]],
                               fill=_solid_color(cell.fill, "#ffffff"), outline="#E9E5DB")
                _draw_text_frame(draw, cell.text_frame,
                                 (tx + 4, ty + 2, col_w[ci] - 8, row_h[ri] - 4))
                tx += col_w[ci]
            ty += row_h[ri]
        return
    if st in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.FREEFORM):
        fillc = _solid_color(sh.fill) if st != MSO_SHAPE_TYPE.TEXT_BOX else None
        linec = _line_color(sh)
        if fillc or linec:
            shape_name = ""
            try:
                shape_name = str(sh.auto_shape_type or "")
            except Exception:
                pass
            if "OVAL" in shape_name:
                draw.ellipse([x, y, x + w, y + hgt], fill=fillc, outline=linec)
            elif "ROUND" in shape_name:
                draw.rounded_rectangle([x, y, x + w, y + hgt],
                                       radius=max(2, min(w, hgt) // 4), fill=fillc, outline=linec)
            else:
                draw.rectangle([x, y, x + w, y + hgt], fill=fillc, outline=linec)
        if sh.has_text_frame:
            _draw_text_frame(draw, sh.text_frame, (x, y, w, hgt))
        return
    if st == MSO_SHAPE_TYPE.LINE or "CONNECTOR" in str(st):
        linec = _line_color(sh, "#999999")
        try:
            draw.line([_px(sh.begin_x), _px(sh.begin_y), _px(sh.end_x), _px(sh.end_y)],
                      fill=linec, width=2)
        except Exception:
            draw.line([x, y, x + w, y + hgt], fill=linec, width=2)
        return
    if sh.has_text_frame:
        _draw_text_frame(draw, sh.text_frame, (x, y, w, hgt))


def render_first_slide(data: bytes, width: int = PREVIEW_WIDTH) -> bytes | None:
    """The title slide of a .pptx as PNG bytes (`width` px wide), or None when the bytes
    don't open as a deck — the caller keeps its extension-badge fallback (graceful, W6)."""
    try:
        from PIL import Image, ImageDraw
        from pptx import Presentation
    except ImportError:
        return None
    try:
        prs = Presentation(io.BytesIO(data))
        slide = next(iter(prs.slides))
    except Exception:
        return None
    full_w, full_h = _px(prs.slide_width), _px(prs.slide_height)
    if full_w <= 0 or full_h <= 0:
        return None
    bg = "#FFFFFF"
    try:
        bg = _solid_color(slide.background.fill, "#FFFFFF")
    except Exception:
        pass
    img = Image.new("RGB", (full_w, full_h), bg)
    draw = ImageDraw.Draw(img)
    for sh in slide.shapes:
        try:
            _render_shape(img, draw, sh)
        except Exception:
            continue                       # one broken shape never sinks the preview
    if width and full_w > width:
        img = img.resize((width, max(1, round(full_h * width / full_w))), Image.LANCZOS)
    out = io.BytesIO()
    img.save(out, format="PNG")
    return out.getvalue()
