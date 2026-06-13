"""Native PowerPoint rendering for reports — turns a neutral slide model into a branded .pptx.

This layer is domain-agnostic (no Sonaloop models, no i18n): the report service assembles a list of
plain slide dicts and hands them here. Charts are NATIVE python-pptx charts (editable in PowerPoint),
coloured from the Sonaloop brand series palette so a deck reads on-brand. Lazy-imported by the caller
so the package degrades gracefully when python-pptx is absent.

Slide model (list of dicts):
  {"kind": "title", "title": str, "subtitle": str, "lead": str}
  {"kind": "content", "heading": str,
     "bullets": [(level:int, text:str)],          # level 0 = paragraph, 1+ = nested bullet
     "chart": {...} | None,                         # see below
     "image": "/abs/path.png" | None,               # a figure image (prototype shot / avatar / asset)
     "footnote": str}
Master-template kinds (the deck taxonomy — single source: sonaloop-design/deck.data.mjs,
vendored as sonaloop/_deck.py; every layout is previewed at #/deck in the design docs and
_deck.SAMPLE_SLIDES carries a placeholder example of each):
  {"kind": "cover", "eyebrow", "title", "subtitle", "meta", "date"}
  {"kind": "agenda", "heading", "items": [str]}
  {"kind": "section", "num", "title", "subtitle"}
  {"kind": "summary", "heading", "items": [{"title", "text"}]}
  {"kind": "insight"|"recommendation"|"risk", "tone", "num", "statement",
     "support": [str], "chart": {...}|None, "meta", "footnote"}
  {"kind": "quote", "text", "attribution", "role"}
  {"kind": "voices", "heading", "items": [{"name", "role", "sentiment", "text"}]}
  {"kind": "stats", "heading", "items": [{"label", "value", "sub"}]}
  {"kind": "chart", "num", "heading", "chart": {...}, "footnote"}
  {"kind": "comparison", "heading", "left": {"title", "items"}, "right": {"title", "items"}}
  {"kind": "timeline", "heading", "steps": [{"label", "title", "text"}]}
  {"kind": "closing", "title", "text", "meta", "contact"}
Chart shapes (one per design-system chart `of`):
  {"type": "bar",  "categories": [str], "values": [num]}
  {"type": "pie",  "categories": [str], "values": [num]}
  {"type": "stacked_bar", "rows": [{"label": str, "segments": [{"label": str, "value": num}]}]}
  {"type": "diverging_bar", "rows": [{"label": str, "positive": num, "negative": num}],
     "positive_label": str, "negative_label": str}
  {"type": "gauge", "items": [{"label": str, "value": num, "max": num}]}
  {"type": "dot_plot", "rows": [{"label": str, "values": [num]}], "min": num, "max": num, "unit": str}
  {"type": "heatmap", "columns": [str], "rows": [{"label": str, "values": [num]}]}
  {"type": "line", "series": [{"label": str, "points": [num]}], "labels": [str], "target": num}
  {"type": "stacked_area", "series": [{"label": str, "points": [num]}], "labels": [str]}
  {"type": "column", "items": [{"label": str, "value": num} | {"label": str, "segments": […]}]}
  {"type": "progress_strip", "items": [{"label": str, "value": num}]}
  {"type": "stats", "items": [{"label": str, "value": num|str, "sub": str}]}
  {"type": "scatter", "points": [{"x":num,"y":num,"label":str}], "x_label": str, "y_label": str}
"""
from __future__ import annotations

import base64
import io
from typing import Any

# Sonaloop brand (light) — from the vendored deck master template (_deck.py, generated out of
# sonaloop-design/deck.data.mjs, which derives them from tokens.data.mjs).
from ._deck import PALETTE as _PALETTE, TONES as _TONES, TYPE as _TYPE

_INK = _PALETTE["ink"]
_MUTED = _PALETTE["muted"]
_FAINT = _PALETTE["faint"]
_ACCENT = _PALETTE["accent"]
_BG = _PALETTE["bg"]
_PANEL = _PALETTE["panel"]
# Series palette (accent · violet · blue · green · amber · red · skep).
_SERIES = list(_PALETTE["series"])
_LINE = _PALETTE["line"]
_SURFACE2 = _PALETTE["surface2"]
_GREEN, _AMBER, _RED = _PALETTE["green"], _PALETTE["amber"], _PALETTE["red"]
_ACCENT_WEAK = _PALETTE["accentWeak"]
# Role-based type sizes (pt) from the master template.
_TS = {k: v.get("size", 13) for k, v in _TYPE.items()}


# Number formatting + the chart painters live in _pptx_charts (split for the LOC bar).
from . import _pptx_charts as _pc
from ._pptx_charts import _num


def render(slides: list[dict], *, title: str = "Report") -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = title
    prs.core_properties.author = "Sonaloop"
    blank = prs.slide_layouts[6]
    W, H = prs.slide_width, prs.slide_height
    rgb = lambda hexv: RGBColor.from_string(hexv)

    def _bg(slide, hexv=_BG):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(hexv)

    def _box(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tb.text_frame.word_wrap = True
        return tb.text_frame

    def _run(p, text, *, size=14, bold=False, color=_INK, italic=False):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = "Geist"
        r.font.color.rgb = rgb(color)
        return r

    def _rule(slide, l, t, w):
        bar = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Pt(3.5))
        bar.fill.solid(); bar.fill.fore_color.rgb = rgb(_ACCENT); bar.line.fill.background()
        _noshadow(bar)

    def _footer(slide):
        ft = _box(slide, W - Inches(5.0), H - Inches(0.42), Inches(4.3), Inches(0.3))
        p = ft.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        _run(p, title, size=9, color=_FAINT)

    # ── title slide ──────────────────────────────────────────────────────────
    def _title_slide(s):
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        _rule(slide, 0.92, 2.0, 1.0)
        tf = _box(slide, Inches(0.9), Inches(2.15), W - Inches(1.8), Inches(3.4))
        p0 = tf.paragraphs[0]
        eb = _run(p0, (s.get("eyebrow", "") or "").upper(), size=12, bold=True, color=_ACCENT)
        eb.font.name = "Geist Mono"
        p1 = tf.add_paragraph(); p1.space_before = Pt(8)
        _run(p1, s.get("title", title), size=38, bold=True)
        if s.get("subtitle"):
            p2 = tf.add_paragraph(); p2.space_before = Pt(12)
            _run(p2, s["subtitle"], size=14, color=_MUTED)
        if s.get("lead"):
            p3 = tf.add_paragraph(); p3.space_before = Pt(20)
            _run(p3, s["lead"], size=17, color=_INK)
        _footer(slide)

    # ── content slide ────────────────────────────────────────────────────────
    _CALLOUT_RGB = {"accent": _ACCENT, "green": _GREEN, "amber": _AMBER}

    def _content_slide(s):
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        # heading: mono section number + bold title, accent rule beneath
        htf = _box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
        hp = htf.paragraphs[0]
        if s.get("num"):
            r = _run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT); r.font.name = "Geist Mono"
        _run(hp, s.get("heading", ""), size=24, bold=True)
        _rule(slide, 0.72, 1.34, 0.85)

        has_visual = bool(s.get("chart") or s.get("image"))
        # prose caps at ~10in without a visual (full-frame lines read like terminal output).
        body_w = 6.4 if has_visual else min(W.inches - 1.4, 10.0)
        region_top = 1.62
        region_bot = H.inches - (0.78 if s.get("footnote") else 0.5)
        region_h = region_bot - region_top

        # Segment the blocks: runs of flowing text (p/li/h/quote) interleaved with boxed callouts.
        _flow_sz = {"li": 13, "quote": 15, "h": 15}
        segs = []; run_blocks = []
        def _flush_run():
            if run_blocks:
                h = sum(_est_h(b.get("text", ""), body_w - 0.2, _flow_sz.get(b.get("type", "p"), 14)) + 0.10
                        for b in run_blocks)
                segs.append(("text", list(run_blocks), h)); run_blocks.clear()
        for b in (s.get("blocks") or []):
            if b.get("type") == "callout":
                _flush_run()
                bh = max(0.66, 0.4 + (0.26 if b.get("label") else 0) + _est_h(b.get("text", ""), body_w - 0.5, 13))
                segs.append(("callout", b, bh))
            else:
                run_blocks.append(b)
        _flush_run()

        gap = 0.16
        total = sum(h for _, _, h in segs) + gap * max(0, len(segs) - 1)
        y = region_top + (max(0.0, (region_h - total) / 2) if total < region_h else 0.0)
        for kind_, payload, h in segs:
            if kind_ == "callout":
                _callout_box(slide, 0.7, y, body_w, payload)
            else:
                tf = _box(slide, Inches(0.7), Inches(y), Inches(body_w), Inches(h + 0.4))
                first = True
                for b in payload:
                    p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
                    bt = b.get("type", "p"); txt = b.get("text", "")
                    if bt == "li":
                        p.space_after = Pt(5)
                        _run(p, "•  ", size=13, bold=True, color=_ACCENT); _run(p, txt, size=13, color=_INK)
                    elif bt == "quote":
                        p.space_before = Pt(4); p.space_after = Pt(6); _run(p, txt, size=15, italic=True, color=_MUTED)
                    elif bt == "h":
                        p.space_before = Pt(8); p.space_after = Pt(2); _run(p, txt, size=15, bold=True, color=_INK)
                    else:
                        p.space_after = Pt(6); _run(p, txt, size=14, color=_INK)
            y += h + gap

        rx = 7.35; rw = W.inches - rx - 0.7
        if s.get("chart"):
            _chart(slide, s["chart"], Inches(rx), Inches(region_top), Inches(rw), Inches(region_h))
        elif s.get("image"):
            try:
                pic = slide.shapes.add_picture(s["image"], Inches(rx), Inches(region_top))
                sc = min(Inches(rw) / pic.width, Inches(region_h) / pic.height)
                pic.width = int(pic.width * sc); pic.height = int(pic.height * sc)
                pic.left = Inches(rx) + (Inches(rw) - pic.width) // 2
                pic.top = Inches(region_top) + (Inches(region_h) - pic.height) // 2
                pic.line.color.rgb = rgb(_LINE); pic.line.width = Pt(0.75)
            except Exception:
                pass
        if s.get("footnote"):
            ftf = _box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.5))
            _run(ftf.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
        _footer(slide)

    # ── image slide (prototype screenshots / images / avatars) — fitted + centred ──
    def _image_slide(s):
        slide = prs.slides.add_slide(blank)
        _bg(slide)
        htf = _box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
        hp = htf.paragraphs[0]
        if s.get("num"):
            r = _run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT); r.font.name = "Geist Mono"
        _run(hp, s.get("heading", ""), size=24, bold=True)
        _rule(slide, 0.72, 1.34, 0.85)
        L, T = 0.7, 1.7
        maxw, maxh = W.inches - 1.4, H.inches - T - 0.95
        placed = False
        if s.get("image"):
            try:
                pic = slide.shapes.add_picture(s["image"], Inches(L), Inches(T))
                scale = min(Inches(maxw) / pic.width, Inches(maxh) / pic.height)
                pic.width = int(pic.width * scale); pic.height = int(pic.height * scale)
                pic.left = Inches(L) + (Inches(maxw) - pic.width) // 2
                pic.top = Inches(T) + (Inches(maxh) - pic.height) // 2
                pic.line.color.rgb = rgb(_LINE); pic.line.width = Pt(0.75)
                if s.get("caption"):
                    cap_t = (pic.top + pic.height) / 914400 + 0.06
                    _text(slide, L, cap_t, maxw, 0.3, s["caption"], size=10, color=_MUTED, align=PP_ALIGN.CENTER)
                placed = True
            except Exception:
                pass
        if not placed:
            # missing/unloadable file → a quiet placeholder panel (the master-template docs case)
            ph = maxh - (0.35 if s.get("caption") else 0)
            _rrect(slide, L, T, maxw, ph, _SURFACE2, radius=0.03, line=_LINE)
            _text(slide, L, T + ph / 2 - 0.2, maxw, 0.4, "image — fitted & centred",
                  size=12, color=_FAINT, align=PP_ALIGN.CENTER)
            if s.get("caption"):
                _text(slide, L, T + ph + 0.08, maxw, 0.3, s["caption"], size=10, color=_MUTED, align=PP_ALIGN.CENTER)
        _footer(slide)

    from pptx.oxml.ns import qn
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
    from pptx.enum.text import MSO_ANCHOR

    # ── chart primitives: native, EDITABLE shapes, pixel-matched to the .sl-chart components ──
    def _noshadow(sp):
        # autoshapes/connectors carry a <p:style> with a theme effectRef (drop shadow). The DS shapes
        # are flat — drop the style entirely (explicit fill/line are set on every shape anyway).
        try:
            sp.shadow.inherit = False
            st = sp._element.find(qn("p:style"))
            if st is not None:
                sp._element.remove(st)
        except Exception:
            pass

    def _text(slide, l, t, w, h, text, *, size=11, color=_INK, bold=False,
              anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.LEFT, rot=0):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        if rot:
            tb.rotation = rot
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        p = tf.paragraphs[0]; p.alignment = align
        _run(p, text, size=size, color=color, bold=bold)
        return tb

    def _rrect(slide, l, t, w, h, color, *, radius=0.5, line=None):
        sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(color)
        if line:
            sp.line.color.rgb = rgb(line); sp.line.width = Pt(1)
        else:
            sp.line.fill.background()
        _noshadow(sp)
        return sp

    def _connector(slide, x1, y1, x2, y2, color, *, width=0.75, dash=False):
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        ln.line.color.rgb = rgb(color); ln.line.width = Pt(width)
        if dash:
            try:
                ln.line._get_or_add_ln().append(ln.line._get_or_add_ln().makeelement(qn("a:prstDash"), {"val": "dash"}))
            except Exception:
                pass
        _noshadow(ln)
        return ln

    def _dot(slide, cxp, cyp, d, fill, edge, num):
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cxp - d / 2), Inches(cyp - d / 2), Inches(d), Inches(d))
        ov.fill.solid(); ov.fill.fore_color.rgb = rgb(fill)
        ov.line.color.rgb = rgb(edge); ov.line.width = Pt(1.5)
        _noshadow(ov)
        tf = ov.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        _run(p, str(num), size=9, bold=True, color=edge)
        return ov

    # ── presentation-layout helpers: vertical balance + boxed callouts ───────────
    def _est_h(text, width_in, size, *, lf=1.34):
        """Rough rendered height (inches) of wrapped text — used to vertically balance a slide."""
        if not text:
            return size * lf / 72.0
        cpl = max(8, int(width_in / (size * 0.0072)))      # chars per line at this width/size
        line = 1; cur = 0
        for word in str(text).split():
            wl = len(word) + 1
            if cur + wl > cpl and cur:
                line += 1; cur = wl
            else:
                cur += wl
        return line * size * lf / 72.0

    _CALLOUT_BAR = {"accent": _ACCENT, "green": _GREEN, "amber": _AMBER,
                    "blue": _PALETTE.get("blue", _ACCENT), "red": _RED}
    _CALLOUT_TINT = {"accent": _ACCENT_WEAK, "green": "E7F3EC", "amber": "F6ECDD",
                     "blue": "E4EEF7", "red": "F7E6E9"}

    def _callout_box(slide, x, y, w, block, *, size=13):
        """A boxed callout: tinted rounded panel + colour bar + optional icon + label + body. Returns height."""
        kind = block.get("kind") or "accent"
        bar = _CALLOUT_BAR.get(kind, _ACCENT); tint = _CALLOUT_TINT.get(kind, _ACCENT_WEAK)
        label = (block.get("label") or "").strip(); txt = block.get("text") or ""
        icon_b64 = _da.ICONS.get(block.get("icon") or "", {}).get("accent")
        pad = 0.2; inner_w = w - 0.5 - (0.72 if icon_b64 else 0)
        h = max(0.66, pad * 2 + (0.26 if label else 0) + _est_h(txt, inner_w, size))
        _rrect(slide, x, y, w, h, tint, radius=0.06)
        _rrect(slide, x, y + 0.12, 0.07, h - 0.24, bar, radius=0.5)
        tx = x + 0.3
        if icon_b64:
            isz = min(0.5, h - 0.3); _pic(slide, icon_b64, x + 0.28, y + (h - isz) / 2, isz, isz); tx = x + 0.28 + isz + 0.18
        tf = _box(slide, Inches(tx), Inches(y + pad - 0.04), Inches(x + w - 0.24 - tx), Inches(h - pad))
        if label:
            _mono_run(_run(tf.paragraphs[0], label.upper(), size=10, bold=True, color=bar))
            bp = tf.add_paragraph(); bp.space_before = Pt(3)
            _run(bp, txt, size=size, color=_INK)
        else:
            _run(tf.paragraphs[0], txt, size=size, color=_INK)
        return h

    # The chart painters (bar/pie/…/scatter) live in _pptx_charts; they draw through the
    # SAME shape primitives via this ctx, so slide and chart layers can't drift apart.
    from types import SimpleNamespace
    _ctx = SimpleNamespace(text=_text, rrect=_rrect, connector=_connector, dot=_dot,
                           run=_run, noshadow=_noshadow)

    def _chart(slide, ch, x, y, cx, cy):
        _pc.draw(_ctx, slide, ch, x, y, cx, cy)

    # ── brand assets — vendored _deck_assets (icons/logos rasterized + canvases recompressed
    # at design time by sonaloop-design/scripts/gen-deck.mjs; PPTX can't embed SVG). Unknown
    # asset names degrade to the unbranded layout. ─────────────────────────────────────
    from . import _deck_assets as _da

    def _pic(slide, b64, l, t, w, h):
        pic = slide.shapes.add_picture(io.BytesIO(base64.b64decode(b64)),
                                       Inches(l), Inches(t), Inches(w), Inches(h))
        _noshadow(pic)
        return pic

    def _pic_cover(slide, b64, l, t, w, h):
        """Aspect-preserving fill (CSS `background: cover`): crop the source to the box ratio."""
        pic = _pic(slide, b64, l, t, w, h)
        try:
            iw, ih = pic.image.size
            sa, ta = iw / ih, w / h
            if sa > ta:
                pic.crop_left = pic.crop_right = (1 - ta / sa) / 2
            elif sa < ta:
                pic.crop_top = pic.crop_bottom = (1 - sa / ta) / 2
        except Exception:
            pass
        return pic

    def _icon_chip(slide, x, y, size, name, *, bg=_ACCENT_WEAK):
        """A hi-fi icon in a tinted rounded chip (the pillars treatment, reusable on any card).
        Returns True if the icon exists in the embedded deck set, else draws nothing."""
        b64 = _da.ICONS.get(name or "", {}).get("accent")
        if not b64:
            return False
        _rrect(slide, x, y, size, size, bg, radius=0.26)
        pad = size * 0.24
        _pic(slide, b64, x + pad, y + pad, size - 2 * pad, size - 2 * pad)
        return True

    def _logo_row(slide, x, y, mark=0.42):
        """The brand moment: mark + wordmark ("sona" ink · "loop" muted)."""
        b64 = _da.LOGOS.get("sonaloop")
        if b64:
            _pic(slide, b64, x, y + 0.04, mark, mark)
        tf = _box(slide, Inches(x + mark + 0.14), Inches(y), Inches(2.6), Inches(mark + 0.08))
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        _run(p, "sona", size=16, bold=True)
        _run(p, "loop", size=16, bold=True, color=_MUTED)

    # ── master-template painters — geometry mirrors sonaloop-design site/deck.preview.mjs
    # painter-for-painter (single source: deck.data.mjs / vendored _deck.py), so the docs
    # previews at #/deck and the exported deck look the same. ──────────────────────────
    def _mono_run(r):
        r.font.name = "Geist Mono"
        return r

    def _heading_band(slide, s, default=""):
        htf = _box(slide, Inches(0.7), Inches(0.5), W - Inches(1.4), Inches(0.9))
        hp = htf.paragraphs[0]
        if s.get("num"):
            _mono_run(_run(hp, s["num"] + "   ", size=16, bold=True, color=_FAINT))
        _run(hp, s.get("heading", default) or default, size=_TS["title"], bold=True)
        _rule(slide, 0.72, 1.34, 0.85)

    def _grid_cells(items):
        """2-column card grid under the heading band → [(x, y, w, h), …] in inches."""
        rows = max(1, (len(items) + 1) // 2)
        gap = 0.25
        cw = (W.inches - 1.4 - gap) / 2
        ch = (H.inches - 1.75 - 0.85 - gap * (rows - 1)) / rows
        return [(0.7 + (i % 2) * (cw + gap), 1.75 + (i // 2) * (ch + gap), cw, ch)
                for i in range(len(items))]

    def _oval(slide, l, t, d, fill):
        ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
        ov.fill.solid(); ov.fill.fore_color.rgb = rgb(fill); ov.line.fill.background()
        _noshadow(ov)
        return ov

    def _initials_chip(slide, l, t, d, name):
        ov = _oval(slide, l, t, d, _ACCENT_WEAK)
        tf = ov.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pq = tf.paragraphs[0]; pq.alignment = PP_ALIGN.CENTER
        ini = "".join(w[0] for w in str(name or "?").split()[:2]).upper() or "?"
        _run(pq, ini, size=10, bold=True, color=_ACCENT)

    def _cover_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        canvas = _da.CANVASES.get(s.get("canvas") or "")
        text_w = W.inches - (4.2 + 1.3 if canvas else 1.8)
        if canvas:
            _pic_cover(slide, canvas, W.inches - 4.2, 0, 4.2, H.inches)
        if s.get("logo"):
            _logo_row(slide, 0.9, 0.55)
        _rule(slide, 0.92, 2.0, 1.0)
        tf = _box(slide, Inches(0.9), Inches(2.15), Inches(text_w), Inches(3.4))
        _mono_run(_run(tf.paragraphs[0], (s.get("eyebrow", "") or "").upper(),
                       size=_TS["eyebrow"], bold=True, color=_ACCENT))
        p1 = tf.add_paragraph(); p1.space_before = Pt(10)
        _run(p1, s.get("title", title), size=_TS["display"], bold=True)
        if s.get("subtitle"):
            p2 = tf.add_paragraph(); p2.space_before = Pt(12)
            _run(p2, s["subtitle"], size=16, color=_MUTED)
        if s.get("meta"):
            mt = _box(slide, Inches(0.9), H - Inches(1.0), Inches(text_w - 2.6), Inches(0.4))
            _mono_run(_run(mt.paragraphs[0], s["meta"], size=11, color=_FAINT))
        if s.get("date"):
            dt = _box(slide, Inches(0.9 + text_w - 2.5), H - Inches(1.0), Inches(2.5), Inches(0.4))
            pd = dt.paragraphs[0]; pd.alignment = PP_ALIGN.RIGHT
            _run(pd, s["date"], size=11, color=_MUTED)

    def _agenda_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s, "Contents")
        for i, item in enumerate(s.get("items") or []):
            y = 1.95 + i * 0.62
            nt = _text(slide, 0.9, y, 0.55, 0.5, str(i + 1).zfill(2), size=14, bold=True, color=_ACCENT)
            _mono_run(nt.text_frame.paragraphs[0].runs[0])
            _text(slide, 1.5, y, W.inches - 2.6, 0.5, str(item), size=16)
        _footer(slide)

    def _section_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        bt = _text(slide, 0.85, 1.0, W.inches - 1.8, 2.6, s.get("num", ""),
                   size=_TS["bignum"], bold=True, color=_ACCENT_WEAK, anchor=MSO_ANCHOR.TOP)
        if s.get("num"):
            _mono_run(bt.text_frame.paragraphs[0].runs[0])
        _rule(slide, 0.92, 3.95, 0.85)
        tf = _box(slide, Inches(0.9), Inches(4.15), W - Inches(1.8), Inches(2.2))
        _run(tf.paragraphs[0], s.get("title", ""), size=32, bold=True)
        if s.get("subtitle"):
            p2 = tf.add_paragraph(); p2.space_before = Pt(10)
            _run(p2, s["subtitle"], size=_TS["subtitle"], color=_MUTED)
        _footer(slide)

    def _canvas_section_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        canvas = _da.CANVASES.get(s.get("canvas") or "")
        if canvas:
            _pic_cover(slide, canvas, 0, 0, W.inches, H.inches)
        y0 = H.inches - 2.65
        _rrect(slide, 0.9, y0, 6.4, 1.75, _PANEL, radius=0.09)
        # the card content flows (mirrors the preview): num is optional, title/subtitle move up
        ty = y0 + 0.24
        if s.get("num"):
            nt = _text(slide, 1.24, y0 + 0.22, 5.7, 0.28, s["num"], size=13, bold=True,
                       color=_ACCENT, anchor=MSO_ANCHOR.TOP)
            _mono_run(nt.text_frame.paragraphs[0].runs[0])
            ty = y0 + 0.52
        _text(slide, 1.24, ty, 5.7, 0.55, s.get("title", ""), size=_TS["title"],
              bold=True, anchor=MSO_ANCHOR.TOP)
        if s.get("subtitle"):
            _text(slide, 1.24, ty + 0.6, 5.7, 0.5, s["subtitle"], size=12, color=_MUTED,
                  anchor=MSO_ANCHOR.TOP)

    def _pillars_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        items = s.get("items") or []
        n = max(len(items), 1)
        gap = 0.3
        cw = (W.inches - 1.4 - gap * (n - 1)) / n
        cy = 2.0
        for i, it in enumerate(items):
            x = 0.7 + i * (cw + gap)
            _rrect(slide, x, cy, 0.72, 0.72, _ACCENT_WEAK, radius=0.19)
            icon = _da.ICONS.get(it.get("icon") or "", {}).get("accent")
            if icon:
                _pic(slide, icon, x + 0.13, cy + 0.13, 0.46, 0.46)
            _text(slide, x, cy + 0.9, cw, 0.4, it.get("title", ""), size=14, bold=True,
                  anchor=MSO_ANCHOR.TOP)
            _text(slide, x, cy + 1.32, cw, H.inches - cy - 2.3, it.get("text", ""), size=11.5,
                  color=_MUTED, anchor=MSO_ANCHOR.TOP)
        _footer(slide)

    def _summary_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s, "Executive summary")
        items = s.get("items") or []
        for it, (cx, cy, cw, ch) in zip(items, _grid_cells(items)):
            _rrect(slide, cx, cy, cw, ch, _PANEL, radius=0.05, line=_LINE)
            if it.get("icon") and _icon_chip(slide, cx + 0.28, cy + 0.26, 0.52, it["icon"]):
                tx0 = cx + 0.96
            else:
                _rrect(slide, cx + 0.26, cy + 0.27, 0.05, 0.21, _ACCENT, radius=0.3); tx0 = cx + 0.42
            # ONE flowing frame (title → body → meta): a wrapping title can never overlap the body
            tf = _box(slide, Inches(tx0), Inches(cy + 0.18), Inches(cx + cw - 0.26 - tx0), Inches(ch - 0.36))
            _run(tf.paragraphs[0], it.get("title", ""), size=15, bold=True)
            if it.get("text"):
                pb = tf.add_paragraph(); pb.space_before = Pt(6)
                _run(pb, it["text"], size=12, color=_MUTED)
            if it.get("meta"):     # quiet card meta (e.g. effort·value scores), mono + faint
                pm = tf.add_paragraph(); pm.space_before = Pt(6)
                _mono_run(_run(pm, str(it["meta"]), size=9, color=_FAINT))
        _footer(slide)

    def _insight_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        tone = _TONES.get(s.get("tone") or s.get("kind") or "insight") or _TONES["insight"]
        tc = _PALETTE.get(tone.get("color", "accent"), _ACCENT)
        et = _box(slide, Inches(0.7), Inches(0.55), W - Inches(1.4), Inches(0.4))
        ep = et.paragraphs[0]
        _mono_run(_run(ep, str(s.get("eyebrow") or tone.get("label", "Insight")).upper(),
                       size=_TS["eyebrow"], bold=True, color=tc))
        if s.get("num"):
            _mono_run(_run(ep, "  ·  " + s["num"], size=_TS["eyebrow"], bold=True, color=_FAINT))
        has_chart = bool(s.get("chart"))
        body_w = 6.9 if has_chart else W.inches - 1.4
        stmt_size, sup_size = 30, 15
        # Content region between the eyebrow and the footnote band; balance the block in it.
        region_top = 1.4
        region_bot = H.inches - (0.95 if (s.get("meta") or s.get("footnote")) else 0.6)
        region_h = region_bot - region_top
        # full-height accent bar (always reads intentional, no matter the block height)
        _rrect(slide, 0.7, region_top, 0.055, region_h, tc, radius=0.3)
        support = [str(t) for t in (s.get("support") or [])]
        est = _est_h(s.get("statement", ""), body_w - 0.3, stmt_size) + 0.18
        est += sum(_est_h(t, body_w - 0.5, sup_size) + 0.11 for t in support)
        anchor = MSO_ANCHOR.MIDDLE if est < region_h - 0.2 else MSO_ANCHOR.TOP
        cf = _box(slide, Inches(0.98), Inches(region_top), Inches(body_w - 0.28), Inches(region_h))
        cf.vertical_anchor = anchor
        _run(cf.paragraphs[0], s.get("statement", ""), size=stmt_size, bold=True)
        for i, t in enumerate(support):
            sp = cf.add_paragraph()
            sp.space_before = Pt(16 if i == 0 else 7); sp.space_after = Pt(0)
            _run(sp, "•  ", size=sup_size, bold=True, color=tc)
            _run(sp, t, size=sup_size)
        if has_chart:
            _chart(slide, s["chart"], Inches(7.55), Inches(1.55), Inches(5.05), Inches(4.5))
        if s.get("meta"):
            mt = _box(slide, Inches(0.7), H - Inches(0.85), Inches(7.0), Inches(0.3))
            _mono_run(_run(mt.paragraphs[0], s["meta"], size=11, bold=True, color=tc))
        if s.get("footnote"):
            ft = _box(slide, Inches(0.7), H - Inches(0.6), W - Inches(1.4), Inches(0.3))
            _run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
        _footer(slide)

    def _quote_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _text(slide, 1.35, 1.05, 1.6, 1.4, "“", size=120, bold=True, color=_ACCENT_WEAK,
              anchor=MSO_ANCHOR.TOP)
        # large, vertically-centred quote — the emotional anchor of the chapter
        _text(slide, 1.9, 1.65, W.inches - 3.4, 3.75, s.get("text", ""), size=30,
              anchor=MSO_ANCHOR.MIDDLE)
        _initials_chip(slide, 1.9, 5.95, 0.36, s.get("attribution"))
        nt = _text(slide, 2.42, 5.97, W.inches - 4.3, 0.4, s.get("attribution", ""), size=13, bold=True)
        _run(nt.text_frame.paragraphs[0], "   " + s.get("role", ""), size=11, color=_MUTED)
        _footer(slide)

    # canonical stance terms → tone colour (stance_scale.json roles; legacy "opposed" kept).
    _SENTIMENT = {"support": _GREEN, "conditional": _AMBER, "neutral": _MUTED,
                  "skeptical": _AMBER, "oppose": _RED, "opposed": _RED}

    def _voices_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s, "Voices")
        items = s.get("items") or []
        if len(items) <= 2:                       # 1–2 voices: full-width stacked cards (quote-scale)
            gap = 0.25
            chh = (H.inches - 1.75 - 0.85 - gap) / 2
            cells = [(0.7, 1.75 + i * (chh + gap), W.inches - 1.4, chh) for i in range(len(items))]
        else:
            cells = _grid_cells(items)
        for it, (cx, cy, cw, ch) in zip(items, cells):
            _rrect(slide, cx, cy, cw, ch, _PANEL, radius=0.05, line=_LINE)
            _initials_chip(slide, cx + 0.24, cy + 0.18, 0.3, it.get("name"))
            nt = _text(slide, cx + 0.62, cy + 0.16, cw - 2.3, 0.34, it.get("name", ""), size=12, bold=True)
            _run(nt.text_frame.paragraphs[0], "   " + it.get("role", ""), size=10, color=_MUTED)
            sc = _SENTIMENT.get((it.get("sentiment") or "").lower(), _MUTED)
            st = _text(slide, cx + cw - 1.85, cy + 0.16, 1.6, 0.34,
                       str(it.get("sentiment_label") or it.get("sentiment") or "").upper(),
                       size=9, bold=True, color=sc, align=PP_ALIGN.RIGHT)
            if it.get("sentiment"):
                _mono_run(st.text_frame.paragraphs[0].runs[0])
            bf = _box(slide, Inches(cx + 0.24), Inches(cy + 0.62), Inches(cw - 0.48), Inches(ch - 0.8))
            _run(bf.paragraphs[0], it.get("text", ""), size=12 if len(items) > 2 else 13)
        _footer(slide)

    def _stats_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        items = s.get("items") or []
        n = max(len(items), 1); gap = 0.28
        tw = (W.inches - 1.4 - gap * (n - 1)) / n
        th = 3.0; ty = 1.75 + (H.inches - 1.75 - 0.6 - th) / 2     # centre the KPI band below the heading
        for i, it in enumerate(items):
            tx = 0.7 + i * (tw + gap)
            _rrect(slide, tx, ty, tw, th, _PANEL, radius=0.05, line=_LINE)
            if it.get("icon") and _icon_chip(slide, tx + 0.28, ty + 0.32, 0.66, it["icon"]):
                tf = _box(slide, Inches(tx + 0.28), Inches(ty + 1.18), Inches(tw - 0.52), Inches(th - 1.4))
                tf.vertical_anchor = MSO_ANCHOR.TOP
            else:
                _rrect(slide, tx + 0.28, ty + 0.32, 0.05, 0.34, _ACCENT, radius=0.3)
                tf = _box(slide, Inches(tx + 0.28), Inches(ty + 0.24), Inches(tw - 0.52), Inches(th - 0.48))
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            v = it.get("value")
            _run(tf.paragraphs[0], v if isinstance(v, str) else _num(v or 0), size=40, bold=True)
            pl = tf.add_paragraph(); pl.space_before = Pt(8)
            _run(pl, str(it.get("label", "")), size=12, color=_MUTED)
            if it.get("sub"):
                p2 = tf.add_paragraph(); p2.space_before = Pt(3)
                _run(p2, str(it["sub"]), size=10, color=_FAINT)
        _footer(slide)

    def _chart_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        if s.get("chart"):
            _chart(slide, s["chart"], Inches(0.7), Inches(1.8), Inches(W.inches - 1.4), Inches(H.inches - 2.9))
        if s.get("footnote"):
            ft = _box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
            _run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
        _footer(slide)

    def _charts_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        items = s.get("items") or []
        n = max(len(items), 1)
        gap = 0.4
        cw = (W.inches - 1.4 - gap * (n - 1)) / n
        top = 1.75
        ch = H.inches - top - 0.5 - (1.3 if s.get("footnote") else 1.0)
        for i, it in enumerate(items):
            x = 0.7 + i * (cw + gap)
            _text(slide, x, top, cw, 0.35, it.get("title", ""), size=13, bold=True,
                  anchor=MSO_ANCHOR.TOP)
            if it.get("chart"):
                _chart(slide, it["chart"], Inches(x), Inches(top + 0.5), Inches(cw), Inches(ch))
        if s.get("footnote"):
            ft = _box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
            _run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
        _footer(slide)

    def _table_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        cols = s.get("columns") or []
        rows = [list(r) for r in (s.get("rows") or [])]
        if cols:
            tw = W.inches - 1.4
            nrows = len(rows) + 1
            avail = H.inches - 1.9 - (0.95 if s.get("footnote") else 0.6)
            row_h = max(0.5, min(0.95, avail / nrows))      # grow rows to fill the frame (capped)
            th = row_h * nrows
            gf = slide.shapes.add_table(nrows, len(cols),
                                        Inches(0.7), Inches(1.9), Inches(tw), Inches(th))
            tbl = gf.table
            # kill the theme's banding/header styling; the cells below carry the deck's own
            tbl.first_row = False
            tbl.horz_banding = False
            for i in range(nrows):
                tbl.rows[i].height = Inches(row_h)
            for j, c in enumerate(cols):
                cell = tbl.cell(0, j)
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb(_ACCENT_WEAK)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.text_frame.word_wrap = True
                _run(cell.text_frame.paragraphs[0], str(c), size=12.5, bold=True)
            for i, row in enumerate(rows, start=1):
                for j in range(len(cols)):
                    cell = tbl.cell(i, j)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = rgb(_BG if i % 2 else _PANEL)
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    cell.text_frame.word_wrap = True
                    val = str(row[j]) if j < len(row) else ""
                    _run(cell.text_frame.paragraphs[0], val, size=12.5,
                         bold=(j == 0), color=_INK if j == 0 else _MUTED)
        if s.get("footnote"):
            ft = _box(slide, Inches(0.7), H - Inches(0.72), W - Inches(1.4), Inches(0.4))
            _run(ft.paragraphs[0], s["footnote"], size=10, color=_FAINT, italic=True)
        _footer(slide)

    def _comparison_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s)
        gap = 0.3
        cw = (W.inches - 1.4 - gap) / 2
        cy = 1.75; ch = H.inches - cy - 0.85
        for j, (col, accent) in enumerate(((s.get("left") or {}, False), (s.get("right") or {}, True))):
            cx = 0.7 + j * (cw + gap)
            _rrect(slide, cx, cy, cw, ch, _PANEL if accent else _SURFACE2, radius=0.04,
                   line=_ACCENT if accent else _LINE)
            tf = _box(slide, Inches(cx + 0.34), Inches(cy + 0.3), Inches(cw - 0.68), Inches(ch - 0.6))
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            _run(tf.paragraphs[0], col.get("title", ""), size=15, bold=True,
                 color=_ACCENT if accent else _MUTED)
            tf.paragraphs[0].space_after = Pt(14)
            for t in col.get("items") or []:
                cp = tf.add_paragraph(); cp.space_after = Pt(9)
                _run(cp, "•  ", size=13, bold=True, color=_ACCENT if accent else _FAINT)
                _run(cp, str(t), size=13)
        _footer(slide)

    def _timeline_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        _heading_band(slide, s, "Next steps")
        steps = s.get("steps") or []
        n = max(len(steps), 1)
        x0, x1, ly = 1.0, W.inches - 1.0, 3.1
        _connector(slide, x0, ly, x1, ly, _LINE, width=1.2)
        for i, st in enumerate(steps):
            cx = x0 + (i + 0.5) * ((x1 - x0) / n)
            _oval(slide, cx - 0.1, ly - 0.1, 0.2, _ACCENT)
            if st.get("label"):
                lt = _text(slide, cx - 1.3, ly - 0.55, 2.6, 0.3, str(st["label"]).upper(),
                           size=10, bold=True, color=_ACCENT, align=PP_ALIGN.CENTER)
                _mono_run(lt.text_frame.paragraphs[0].runs[0])
            _text(slide, cx - 1.3, ly + 0.25, 2.6, 0.5, st.get("title", ""), size=13, bold=True,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            _text(slide, cx - 1.3, ly + 0.72, 2.6, 1.4, st.get("text", ""), size=11, color=_MUTED,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
        _footer(slide)

    def _closing_slide(s):
        slide = prs.slides.add_slide(blank); _bg(slide)
        if s.get("logo"):
            _logo_row(slide, 0.9, 1.1)
        _rule(slide, 0.92, 2.35, 1.0)
        tf = _box(slide, Inches(0.9), Inches(2.55), Inches(10.2), Inches(H.inches - 3.6))
        _run(tf.paragraphs[0], s.get("title", ""), size=40, bold=True)
        if s.get("text"):
            p1 = tf.add_paragraph(); p1.space_before = Pt(18)
            _run(p1, s["text"], size=16, color=_INK)
        if s.get("contact"):
            p2 = tf.add_paragraph(); p2.space_before = Pt(20)
            _run(p2, s["contact"], size=14, bold=True, color=_ACCENT)
        if s.get("meta"):
            mt = _box(slide, Inches(0.9), H - Inches(1.0), W - Inches(1.8), Inches(0.4))
            _mono_run(_run(mt.paragraphs[0], s["meta"], size=11, color=_FAINT))

    painters = {"title": _title_slide, "cover": _cover_slide, "agenda": _agenda_slide,
                "section": _section_slide, "canvas-section": _canvas_section_slide,
                "pillars": _pillars_slide, "summary": _summary_slide,
                "insight": _insight_slide, "recommendation": _insight_slide, "risk": _insight_slide,
                "quote": _quote_slide, "voices": _voices_slide, "stats": _stats_slide,
                "chart": _chart_slide, "charts": _charts_slide, "table": _table_slide,
                "comparison": _comparison_slide, "timeline": _timeline_slide,
                "closing": _closing_slide, "image": _image_slide}
    for s in slides:
        painters.get(s.get("kind"), _content_slide)(s)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def available() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except Exception:
        return False
