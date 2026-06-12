"""Native chart painters for the PPTX renderer — the design-system chart catalogue drawn as
EDITABLE python-pptx shapes/charts, pixel-matched to the .sl-chart components.

Split out of sonaloop/_pptx.py to keep both modules under the LOC bar (spec/refactor-plan.md),
behaviour-preserving. The slide painters in _pptx.py call draw() with a small ctx of shared
drawing primitives (text / rrect / connector / dot), so both modules paint with ONE set of
shape helpers. Chart shapes (one per design-system chart `of`) are documented in _pptx.py.
"""
from __future__ import annotations

# Brand palette — same vendored source as _pptx.py (sonaloop-design deck.data.mjs → _deck.py).
from ._deck import PALETTE as _PALETTE

_INK = _PALETTE["ink"]
_MUTED = _PALETTE["muted"]
_FAINT = _PALETTE["faint"]
_ACCENT = _PALETTE["accent"]
_BG = _PALETTE["bg"]
_SERIES = list(_PALETTE["series"])
_LINE = _PALETTE["line"]
_SURFACE2 = _PALETTE["surface2"]
_GREEN, _AMBER, _RED = _PALETTE["green"], _PALETTE["amber"], _PALETTE["red"]


def _num(v: float) -> str:
    """Compact number for labels: 8 not 8.0, 7.5 stays 7.5."""
    try:
        f = float(v)
        return str(int(f)) if f.is_integer() else f"{f:g}"
    except (TypeError, ValueError):
        return str(v)


def _leverage_color(x: float, y: float) -> str:
    """Effort·impact dot tint — mirrors sonaloop-design _charts: high value vs effort → green;
    balanced → accent; costly → amber/red."""
    d = y - x
    return _GREEN if d >= 2 else _ACCENT if d >= 1 else _RED if d <= -1 else _AMBER


def _mix(a: str, b: str, t: float) -> str:
    """Blend hex colours a→b by t∈[0,1] (t=0 → a, t=1 → b). Used for the heatmap tint, mirroring the
    DS `color-mix(in srgb, accent p%, surface-2)` — pass _mix(_SURFACE2, _ACCENT, value_fraction)."""
    t = max(0.0, min(1.0, t))
    ca = (int(a[0:2], 16), int(a[2:4], 16), int(a[4:6], 16))
    cb = (int(b[0:2], 16), int(b[2:4], 16), int(b[4:6], 16))
    return "".join(f"{round(ca[i] + (cb[i] - ca[i]) * t):02X}" for i in range(3))


def _rgb(hexv: str):
    from pptx.dml.color import RGBColor
    return RGBColor.from_string(hexv)


def _set_hole_size(plot, val: str):
    """Set the doughnut hole size on the EXISTING `c:holeSize` element (python-pptx's doughnut
    template already carries one at val=50). Appending a second one — the pre-round-3 code path —
    made the chart part schema-invalid (CT_DoughnutChart allows ONE holeSize), so strict renderers
    (PowerPoint) dropped the whole chart: the H5 "legend with NO donut" empty panel."""
    from pptx.oxml.ns import qn
    el = plot._element
    hole = el.find(qn("c:holeSize"))
    if hole is None:
        hole = el.makeelement(qn("c:holeSize"), {})
        el.append(hole)
    hole.set("val", val)


def _clean_area(chart):
    """Transparent chart area (blend into the slide) + brand default font — for the donut ring,
    which stays a native chart (arcs are impractical as shapes)."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    try:
        chart.font.name = "Geist"; chart.font.size = Pt(10); chart.font.color.rgb = _rgb(_MUTED)
        cs = chart._chartSpace
        spPr = cs.makeelement(qn("c:spPr"), {})
        cs.find(qn("c:chart")).addnext(spPr)
        spPr.append(spPr.makeelement(qn("a:noFill"), {}))
        ln = spPr.makeelement(qn("a:ln"), {}); spPr.append(ln)
        ln.append(ln.makeelement(qn("a:noFill"), {}))
    except Exception:
        pass


def _bar_chart(ctx, slide, ch, bx, by, bw, bh):
    cats = [str(c) for c in ch["categories"]]
    vals = [float(v or 0) for v in ch["values"]]
    cols = ch.get("colors") or _SERIES         # optional semantic palette (e.g. stance roles)
    n = max(len(cats), 1); mx = max(vals + [1]) or 1
    label_w = min(1.5, bw * 0.34); value_w = 0.42
    tx = bx + label_w; tw = max(0.5, bw - label_w - value_w)
    row_h = min(bh / n, 0.5); bar_h = min(0.16, row_h * 0.38)
    y0 = by + (bh - row_h * n) / 2          # centre the bar group vertically (DS rows are tight)
    for i, (c, v) in enumerate(zip(cats, vals)):
        ry = y0 + i * row_h; cy0 = ry + (row_h - bar_h) / 2
        ctx.text(slide, bx, ry, label_w - 0.08, row_h, c, size=11, color=_INK)
        ctx.rrect(slide, tx, cy0, tw, bar_h, _SURFACE2)
        ctx.rrect(slide, tx, cy0, max(bar_h, tw * v / mx), bar_h, cols[i % len(cols)])
        ctx.text(slide, tx + tw + 0.05, ry, value_w, row_h, _num(v), size=10, color=_MUTED)


def _donut_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    cats = [str(c) for c in ch["categories"]]
    vals = [float(v or 0) for v in ch["values"]]; total = sum(vals) or 1
    cols = ch.get("colors") or _SERIES         # optional semantic palette (e.g. stance roles)
    size = min(bw * 0.46, bh, 3.4)
    cd = CategoryChartData(); cd.categories = cats; cd.add_series("", tuple(vals))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(bx), Inches(by + (bh - size) / 2),
                                   Inches(size), Inches(size), cd).chart
    chart.has_title = False; chart.has_legend = False
    _clean_area(chart)
    plot = chart.plots[0]; plot.has_data_labels = False
    try:
        for j, pt in enumerate(plot.series[0].points):
            pt.format.fill.solid(); pt.format.fill.fore_color.rgb = _rgb(cols[j % len(cols)])
            pt.format.line.fill.background()
        _set_hole_size(plot, "62")
    except Exception:
        pass
    lx = bx + size + 0.4; lw = max(1.0, bw - size - 0.4)
    rh = min(0.36, bh / max(len(cats), 1)); ly0 = by + (bh - rh * len(cats)) / 2
    for i, (c, v) in enumerate(zip(cats, vals)):
        ry = ly0 + i * rh
        ctx.rrect(slide, lx, ry + rh / 2 - 0.07, 0.14, 0.14, cols[i % len(cols)], radius=0.28)
        tb = ctx.text(slide, lx + 0.26, ry, lw - 0.26, rh, c + "  ", size=11, color=_INK)
        ctx.run(tb.text_frame.paragraphs[0], f"{_num(v)} · {round(v / total * 100)}%", size=10, color=_MUTED)


def _effort_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    pts = ch.get("points") or []
    xlab, ylab = ch.get("x_label", "Effort"), ch.get("y_label", "Value")
    q = (ch.get("quadrants") or ["Quick wins", "Big bets", "Fill-ins", "Time sinks"]) + ["", "", "", ""]
    legend_h = min(bh * 0.46, 0.27 * len(pts) + 0.05)
    ylab_w, xlab_h = 0.32, 0.3
    plot_h = max(1.4, bh - xlab_h - legend_h - 0.2)
    ox, oy, ow, oh = bx + ylab_w, by, bw - ylab_w, plot_h
    # L-shaped frame (left + bottom, like .sl-quad) + dashed centre cross
    ctx.connector(slide, ox, oy, ox, oy + oh, _LINE, width=1)
    ctx.connector(slide, ox, oy + oh, ox + ow, oy + oh, _LINE, width=1)
    ctx.connector(slide, ox + ow / 2, oy, ox + ow / 2, oy + oh, _LINE, dash=True)
    ctx.connector(slide, ox, oy + oh / 2, ox + ow, oy + oh / 2, _LINE, dash=True)
    ctx.text(slide, ox + 0.06, oy + 0.02, ow / 2 - 0.1, 0.22, q[0], size=9, color=_FAINT, anchor=MSO_ANCHOR.TOP)
    ctx.text(slide, ox + ow / 2, oy + 0.02, ow / 2 - 0.06, 0.22, q[1], size=9, color=_FAINT, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.RIGHT)
    ctx.text(slide, ox + 0.06, oy + oh - 0.24, ow / 2 - 0.1, 0.22, q[2], size=9, color=_FAINT, anchor=MSO_ANCHOR.BOTTOM)
    ctx.text(slide, ox + ow / 2, oy + oh - 0.24, ow / 2 - 0.06, 0.22, q[3], size=9, color=_FAINT, anchor=MSO_ANCHOR.BOTTOM, align=PP_ALIGN.RIGHT)
    for i, pt in enumerate(pts, 1):
        px, py = float(pt.get("x", 1)), float(pt.get("y", 1)); col = _leverage_color(px, py)
        ctx.dot(slide, ox + (px - 1) / 4 * ow, oy + (1 - (py - 1) / 4) * oh, 0.3, _BG, col, i)
    ctx.text(slide, ox, by + oh + 0.02, ow, xlab_h, xlab, size=10, color=_INK, align=PP_ALIGN.CENTER)
    ctx.text(slide, bx - 0.5, by + oh / 2 - 0.15, 1.3, 0.3, ylab, size=10, color=_INK, align=PP_ALIGN.CENTER, rot=270)
    # numbered legend (matches the DS dot legend) below the axis label
    leg_y = by + oh + xlab_h + 0.04
    for i, pt in enumerate(pts, 1):
        px, py = float(pt.get("x", 1)), float(pt.get("y", 1)); col = _leverage_color(px, py)
        ry = leg_y + (i - 1) * 0.27
        ctx.dot(slide, bx + 0.11, ry + 0.115, 0.2, _BG, col, i)
        ctx.text(slide, bx + 0.3, ry, bw - 1.0, 0.23, pt.get("label", ""), size=10, color=_INK)
        ctx.text(slide, bx + bw - 0.7, ry, 0.7, 0.23, f"{xlab[:1]}{_num(px)}·{ylab[:1]}{_num(py)}", size=9, color=_MUTED, align=PP_ALIGN.RIGHT)


def _legend_row(ctx, slide, items, lx, ly, lw):
    """A horizontal swatch+label legend (matches the DS .sl-legend--row). items: [(label, color)]."""
    x = lx
    for lbl, col in items:
        if x - lx > lw - 0.6:
            break
        ctx.rrect(slide, x, ly + 0.03, 0.13, 0.13, col, radius=0.28)
        ctx.text(slide, x + 0.19, ly, min(1.7, 0.085 * len(str(lbl)) + 0.2), 0.2, str(lbl), size=9, color=_MUTED)
        x += 0.19 + min(1.7, 0.085 * len(str(lbl)) + 0.2) + 0.18


def _stacked_bar_chart(ctx, slide, ch, bx, by, bw, bh):
    rows = ch.get("rows") or []
    keys = []
    for r in rows:
        for seg in r.get("segments") or []:
            if seg.get("label") not in keys:
                keys.append(seg.get("label"))
    col_of = lambda lbl: _SERIES[(keys.index(lbl) if lbl in keys else 0) % len(_SERIES)]
    totals = [sum(float(s.get("value") or 0) for s in (r.get("segments") or [])) for r in rows]
    n = max(len(rows), 1); mx = max(totals + [1]) or 1
    label_w = min(1.5, bw * 0.30); value_w = 0.42
    tx = bx + label_w; tw = max(0.5, bw - label_w - value_w)
    legend_h = 0.34; avail = bh - legend_h
    row_h = min(avail / n, 0.5); bar_h = min(0.18, row_h * 0.4)
    y0 = by + (avail - row_h * n) / 2
    for i, (r, total) in enumerate(zip(rows, totals)):
        ry = y0 + i * row_h; cy0 = ry + (row_h - bar_h) / 2
        ctx.text(slide, bx, ry, label_w - 0.08, row_h, r.get("label", ""), size=11)
        ctx.rrect(slide, tx, cy0, tw, bar_h, _SURFACE2)
        sx = tx; full = max(bar_h, tw * total / mx)
        for seg in r.get("segments") or []:
            v = float(seg.get("value") or 0)
            if v <= 0:
                continue
            w = full * (v / total) if total else 0
            ctx.rrect(slide, sx, cy0, max(0.02, w), bar_h, col_of(seg.get("label")), radius=0)
            sx += w
        ctx.text(slide, tx + tw + 0.05, ry, value_w, row_h, _num(total), size=10, color=_MUTED)
    _legend_row(ctx, slide, [(k, col_of(k)) for k in keys], tx, by + bh - legend_h + 0.04, tw)


def _diverging_chart(ctx, slide, ch, bx, by, bw, bh):
    rows = ch.get("rows") or []
    pos_l, neg_l = ch.get("positive_label", "Positive"), ch.get("negative_label", "Negative")
    mx = max([max(abs(float(r.get("positive") or 0)), abs(float(r.get("negative") or 0))) for r in rows] + [1]) or 1
    n = max(len(rows), 1)
    label_w = min(1.4, bw * 0.24); value_w = 0.95
    tx = bx + label_w; tw = max(0.6, bw - label_w - value_w); half = tw / 2; cx = tx + half
    legend_h = 0.34; avail = bh - legend_h
    row_h = min(avail / n, 0.5); bar_h = min(0.18, row_h * 0.4)
    y0 = by + (avail - row_h * n) / 2
    for i, r in enumerate(rows):
        ry = y0 + i * row_h; cy0 = ry + (row_h - bar_h) / 2
        pos = max(0.0, float(r.get("positive") or 0)); neg = max(0.0, float(r.get("negative") or 0))
        ctx.text(slide, bx, ry, label_w - 0.06, row_h, r.get("label", ""), size=11)
        ctx.rrect(slide, tx, cy0, tw, bar_h, _SURFACE2)
        negw = half * neg / mx; posw = half * pos / mx
        if negw > 0:
            ctx.rrect(slide, cx - negw, cy0, negw, bar_h, _RED, radius=0)
        if posw > 0:
            ctx.rrect(slide, cx, cy0, posw, bar_h, _GREEN, radius=0)
        ctx.text(slide, tx + tw + 0.05, ry, value_w, row_h, f"+{_num(pos)} · −{_num(neg)}", size=9, color=_MUTED)
    ctx.connector(slide, cx, y0, cx, y0 + row_h * n, _LINE, width=1)
    _legend_row(ctx, slide, [(pos_l, _GREEN), (neg_l, _RED)], tx, by + bh - legend_h + 0.04, tw)


def _gauge_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    items = ch.get("items") or []
    n = max(len(items), 1); gap = 0.3
    size = max(0.8, min((bw - gap * (n - 1)) / n, bh - 0.5, 2.2))
    x0 = bx + (bw - (size * n + gap * (n - 1))) / 2
    for i, it in enumerate(items):
        m = float(it.get("max") or 100) or 1
        v = max(0.0, min(m, float(it.get("value") or 0))); pct = v / m * 100
        gx = x0 + i * (size + gap)
        cd = CategoryChartData(); cd.categories = ["v", "rest"]; cd.add_series("", (v, max(0.0, m - v)))
        chart = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Inches(gx), Inches(by),
                                       Inches(size), Inches(size), cd).chart
        chart.has_title = False; chart.has_legend = False
        _clean_area(chart)
        plot = chart.plots[0]; plot.has_data_labels = False
        try:
            pts = plot.series[0].points
            for idx, col in ((0, _SERIES[i % len(_SERIES)]), (1, _SURFACE2)):
                pts[idx].format.fill.solid(); pts[idx].format.fill.fore_color.rgb = _rgb(col)
                pts[idx].format.line.fill.background()
            _set_hole_size(plot, "70")
        except Exception:
            pass
        ctx.text(slide, gx, by + size / 2 - 0.2, size, 0.4, f"{round(pct)}%", size=16, bold=True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        ctx.text(slide, gx - 0.1, by + size + 0.04, size + 0.2, 0.3, it.get("label", ""), size=10, align=PP_ALIGN.CENTER)


def _dot_plot_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    rows = ch.get("rows") or []
    u = str(ch.get("unit") or "")
    mn = float(ch.get("min", 1)); mx = float(ch.get("max", 5)); span = (mx - mn) or 1
    n = max(len(rows), 1)
    label_w = min(1.6, bw * 0.32); value_w = 0.42; scale_h = 0.3
    tx = bx + label_w; tw = max(0.5, bw - label_w - value_w)
    avail = bh - scale_h; row_h = min(avail / n, 0.5)
    y0 = by + (avail - row_h * n) / 2
    xof = lambda v: tx + max(0.0, min(1.0, (v - mn) / span)) * tw
    for i, r in enumerate(rows):
        ry = y0 + i * row_h; cyc = ry + row_h / 2
        vals = [float(v) for v in (r.get("values") or []) if isinstance(v, (int, float))]
        if not vals:
            continue
        mean = round(sum(vals) / len(vals) * 10) / 10; col = _SERIES[i % len(_SERIES)]
        ctx.text(slide, bx, ry, label_w - 0.08, row_h, r.get("label", ""), size=11)
        ctx.connector(slide, tx, cyc, tx + tw, cyc, _LINE, width=0.75)
        d = 0.13; tint = _mix(col, _BG, 0.45)
        for v in vals:
            ov = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(xof(v) - d / 2), Inches(cyc - d / 2), Inches(d), Inches(d))
            ov.fill.solid(); ov.fill.fore_color.rgb = _rgb(tint); ov.line.fill.background(); ctx.noshadow(ov)
        mh = 0.22
        ctx.rrect(slide, xof(mean) - 0.025, cyc - mh / 2, 0.05, mh, col, radius=0.3)
        ctx.text(slide, tx + tw + 0.05, ry, value_w, row_h, _num(mean) + u, size=10, color=_MUTED)
    ctx.text(slide, tx, by + bh - scale_h + 0.02, 0.6, scale_h, _num(mn) + u, size=9, color=_FAINT)
    ctx.text(slide, tx + tw - 0.6, by + bh - scale_h + 0.02, 0.6, scale_h, _num(mx) + u, size=9, color=_FAINT, align=PP_ALIGN.RIGHT)


def _heatmap_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    cols = [str(c) for c in (ch.get("columns") or [])]
    rows = ch.get("rows") or []
    if not cols or not rows:
        return
    allv = [float(v) for r in rows for v in (r.get("values") or []) if isinstance(v, (int, float))]
    mn = min(allv + [0]); mx = max(allv + [1])
    ncols, nrows = len(cols), len(rows); gap = 0.05
    label_w = min(1.6, bw * 0.28); header_h = 0.3
    cell_w = (bw - label_w - gap * ncols) / ncols
    cell_h = min((bh - header_h - gap * nrows) / max(nrows, 1), 0.6)
    for j, c in enumerate(cols):
        ctx.text(slide, bx + label_w + j * (cell_w + gap), by, cell_w, header_h, c, size=9, color=_MUTED, align=PP_ALIGN.CENTER)
    y0 = by + header_h
    for i, r in enumerate(rows):
        ry = y0 + i * (cell_h + gap)
        ctx.text(slide, bx, ry, label_w - 0.08, cell_h, r.get("label", ""), size=10, anchor=MSO_ANCHOR.MIDDLE)
        vals = r.get("values") or []
        for j in range(ncols):
            cxp = bx + label_w + j * (cell_w + gap)
            v = vals[j] if j < len(vals) and isinstance(vals[j], (int, float)) else None
            if v is None:
                ctx.rrect(slide, cxp, ry, cell_w, cell_h, _SURFACE2, radius=0.12)
            else:
                t = 0.0 if mx == mn else max(0.0, min(1.0, (float(v) - mn) / (mx - mn)))
                ctx.rrect(slide, cxp, ry, cell_w, cell_h, _mix(_SURFACE2, _ACCENT, t), radius=0.12)
                ctx.text(slide, cxp, ry, cell_w, cell_h, _num(v), size=10, color=_INK,
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def _line_chart(ctx, slide, ch, bx, by, bw, bh):
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    lines = [s for s in (ch.get("series") or []) if (s.get("points") or [])]
    if not lines:
        return
    maxlen = max(len(s.get("points") or []) for s in lines)
    labels = ch.get("labels") or []
    cats = [str(x) for x in labels] if len(labels) == maxlen else [str(i + 1) for i in range(maxlen)]
    cd = CategoryChartData(); cd.categories = cats
    for s in lines:
        pts = [float(p) if isinstance(p, (int, float)) else None for p in (s.get("points") or [])]
        cd.add_series(s.get("label", ""), tuple(pts + [None] * (maxlen - len(pts))))
    # Burn-up ideal line — a straight dashed series from the first point up to the target.
    tgt = ch.get("target")
    has_target = isinstance(tgt, (int, float)) and not isinstance(tgt, bool) and maxlen > 1
    if has_target:
        first = next((float(p) for p in (lines[0].get("points") or []) if isinstance(p, (int, float))), 0.0)
        cd.add_series("Target", tuple(first + (float(tgt) - first) * i / (maxlen - 1) for i in range(maxlen)))
    multi = len(lines) > 1 or has_target; legend_h = 0.0
    chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(bx), Inches(by),
                                   Inches(bw), Inches(bh - legend_h), cd).chart
    chart.has_title = False; chart.has_legend = multi
    if multi:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9); chart.legend.font.name = "Geist"; chart.legend.font.color.rgb = _rgb(_MUTED)
    _clean_area(chart)
    for i, ser in enumerate(chart.series):
        is_target = has_target and i == len(lines)
        col = _MUTED if is_target else _SERIES[i % len(_SERIES)]
        try:
            ser.format.line.color.rgb = _rgb(col); ser.format.line.width = Pt(1.25 if is_target else 2)
            ser.smooth = False
            if is_target:
                from pptx.enum.chart import XL_MARKER_STYLE
                from pptx.enum.line import MSO_LINE
                ser.format.line.dash_style = MSO_LINE.DASH
                ser.marker.style = XL_MARKER_STYLE.NONE
        except Exception:
            pass
    try:
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(9); ax.tick_labels.font.name = "Geist"
            ax.tick_labels.font.color.rgb = _rgb(_MUTED)
            ax.format.line.color.rgb = _rgb(_LINE)
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass


def _area_chart(ctx, slide, ch, bx, by, bw, bh):
    """Stacked area / cumulative flow — a native AREA_STACKED chart, series-palette filled."""
    from pptx.util import Inches, Pt
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    bands = [s for s in (ch.get("series") or []) if (s.get("points") or [])]
    if not bands:
        return
    maxlen = max(len(s.get("points") or []) for s in bands)
    labels = ch.get("labels") or []
    cats = [str(x) for x in labels] if len(labels) == maxlen else [str(i + 1) for i in range(maxlen)]
    cd = CategoryChartData(); cd.categories = cats
    for s in bands:
        pts = [float(p) if isinstance(p, (int, float)) else 0.0 for p in (s.get("points") or [])]
        cd.add_series(s.get("label", ""), tuple(pts + [0.0] * (maxlen - len(pts))))
    chart = slide.shapes.add_chart(XL_CHART_TYPE.AREA_STACKED, Inches(bx), Inches(by),
                                   Inches(bw), Inches(bh), cd).chart
    chart.has_title = False; chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM; chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9); chart.legend.font.name = "Geist"; chart.legend.font.color.rgb = _rgb(_MUTED)
    _clean_area(chart)
    for i, ser in enumerate(chart.series):
        col = _SERIES[i % len(_SERIES)]
        try:
            ser.format.fill.solid(); ser.format.fill.fore_color.rgb = _rgb(_mix(col, _BG, 0.35))
            ser.format.line.color.rgb = _rgb(col); ser.format.line.width = Pt(1.25)
        except Exception:
            pass
    try:
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(9); ax.tick_labels.font.name = "Geist"
            ax.tick_labels.font.color.rgb = _rgb(_MUTED)
            ax.format.line.color.rgb = _rgb(_LINE)
        chart.value_axis.has_major_gridlines = False
    except Exception:
        pass


def _column_chart(ctx, slide, ch, bx, by, bw, bh):
    """Vertical bars (the DS column chart) — plain or segment-stacked, value + category labels."""
    from pptx.enum.text import PP_ALIGN
    items = ch.get("items") or []
    keys = []
    for it in items:
        for seg in it.get("segments") or []:
            if seg.get("label") not in keys:
                keys.append(seg.get("label"))
    col_of = lambda lbl: _SERIES[(keys.index(lbl) if lbl in keys else 0) % len(_SERIES)]

    def _total(it):
        if it.get("segments"):
            return sum(float(s.get("value") or 0) for s in it["segments"])
        return float(it.get("value") or 0)

    totals = [_total(it) for it in items]
    if not totals:
        return
    mx = max(totals + [1]) or 1
    n = max(len(items), 1)
    label_h, val_h = 0.26, 0.2
    legend_h = 0.32 if keys else 0.0
    plot_h = max(0.8, bh - label_h - val_h - legend_h)
    gap = min(0.3, bw / n * 0.4)
    cw = (bw - gap * (n - 1)) / n; bar_w = min(cw, 0.55)
    base_y = by + val_h + plot_h
    ctx.connector(slide, bx, base_y, bx + bw, base_y, _LINE, width=1)
    for i, (it, tot) in enumerate(zip(items, totals)):
        cx0 = bx + i * (cw + gap)
        bx0 = cx0 + (cw - bar_w) / 2
        h = plot_h * tot / mx
        if it.get("segments"):
            sy = base_y
            for seg in it["segments"]:
                v = float(seg.get("value") or 0)
                if v <= 0:
                    continue
                sh = h * (v / tot) if tot else 0
                sy -= sh
                ctx.rrect(slide, bx0, sy, bar_w, sh, col_of(seg.get("label")), radius=0)
        elif h > 0:
            ctx.rrect(slide, bx0, base_y - h, bar_w, max(0.02, h), _SERIES[0], radius=0)
        ctx.text(slide, cx0, base_y - h - val_h, cw, val_h, _num(tot), size=9, color=_MUTED, align=PP_ALIGN.CENTER)
        ctx.text(slide, cx0, base_y + 0.03, cw, label_h, str(it.get("label", "")), size=9, color=_MUTED, align=PP_ALIGN.CENTER)
    if keys:
        _legend_row(ctx, slide, [(k, col_of(k)) for k in keys], bx, base_y + label_h + 0.06, bw)


def _progress_strip_chart(ctx, slide, ch, bx, by, bw, bh):
    """One full-width segmented status bar + a count/% legend (the DS progress strip)."""
    items = [it for it in (ch.get("items") or []) if float(it.get("value") or 0) > 0]
    total = sum(float(it["value"]) for it in items)
    if not items or total <= 0:
        return
    bar_h = 0.24; bar_y = by + min(0.4, bh * 0.2)
    ctx.rrect(slide, bx, bar_y, bw, bar_h, _SURFACE2)
    sx = bx
    for i, it in enumerate(items):
        w = bw * float(it["value"]) / total
        ctx.rrect(slide, sx, bar_y, max(0.02, w), bar_h, _SERIES[i % len(_SERIES)], radius=0)
        sx += w
    x = bx; ly = bar_y + bar_h + 0.14
    for i, it in enumerate(items):
        v = float(it["value"])
        lbl = f'{it.get("label", "")}  {_num(v)} · {round(v / total * 100)}%'
        w = min(2.3, 0.07 * len(lbl) + 0.25)
        if x + 0.19 + w > bx + bw:
            x = bx; ly += 0.24
        ctx.rrect(slide, x, ly + 0.03, 0.13, 0.13, _SERIES[i % len(_SERIES)], radius=0.28)
        ctx.text(slide, x + 0.19, ly, w, 0.2, lbl, size=9, color=_MUTED)
        x += 0.19 + w + 0.16


def _stats_chart(ctx, slide, ch, bx, by, bw, bh):
    """KPI number row — label · big value · optional sub, one tile per item."""
    items = ch.get("items") or []
    if not items:
        return
    n = len(items)
    tile_w = bw / n
    for i, it in enumerate(items):
        gx = bx + i * tile_w
        v = it.get("value")
        txt = v if isinstance(v, str) else _num(float(v or 0))
        ctx.text(slide, gx, by + 0.05, tile_w - 0.12, 0.22, str(it.get("label", "")), size=10, color=_MUTED)
        ctx.text(slide, gx, by + 0.28, tile_w - 0.12, 0.5, str(txt), size=24, bold=True)
        if it.get("sub"):
            ctx.text(slide, gx, by + 0.8, tile_w - 0.12, 0.22, str(it["sub"]), size=9, color=_FAINT)


def draw(ctx, slide, ch, x, y, cx, cy):
    """Dispatch one neutral chart dict into its painter. x/y/cx/cy are pptx Length objects."""
    bx, by, bw, bh = x.inches, y.inches, cx.inches, cy.inches
    kind = ch.get("type")
    try:
        if kind == "bar" and ch.get("categories"):
            _bar_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "pie" and ch.get("categories"):
            _donut_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "stacked_bar" and ch.get("rows"):
            _stacked_bar_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "diverging_bar" and ch.get("rows"):
            _diverging_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "gauge" and ch.get("items"):
            _gauge_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "dot_plot" and ch.get("rows"):
            _dot_plot_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "heatmap" and ch.get("columns") and ch.get("rows"):
            _heatmap_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "line" and ch.get("series"):
            _line_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "stacked_area" and ch.get("series"):
            _area_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "column" and ch.get("items"):
            _column_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "progress_strip" and ch.get("items"):
            _progress_strip_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "stats" and ch.get("items"):
            _stats_chart(ctx, slide, ch, bx, by, bw, bh)
        elif kind == "scatter" and ch.get("points"):
            _effort_chart(ctx, slide, ch, bx, by, bw, bh)
    except Exception:
        pass
